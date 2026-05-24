"""
西安建筑科技大学管理学院 工商管理专业学业与职业规划 PPT 生成脚本
使用 python-pptx 生成 .pptx 文件
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ============================================================
# 配色方案
# ============================================================
RED = RGBColor(0xC4, 0x1E, 0x2A)       # 西建大红
DARK = RGBColor(0x33, 0x33, 0x33)       # 深灰
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF5, 0xF5, 0xF5)
MID_GRAY = RGBColor(0x99, 0x99, 0x99)
BLUE = RGBColor(0x1A, 0x73, 0xE8)
GREEN = RGBColor(0x34, 0xA8, 0x53)
ORANGE = RGBColor(0xF5, 0xA6, 0x23)
LIGHT_RED = RGBColor(0xFF, 0xEB, 0xEE)
LIGHT_BLUE = RGBColor(0xE3, 0xF2, 0xFD)
LIGHT_GREEN = RGBColor(0xE8, 0xF5, 0xE9)
LIGHT_ORANGE = RGBColor(0xFF, 0xF3, 0xE0)

OUTPUT_PATH = r"D:\develop\cursor Pro\shao-xiaoli-attendance\工商管理专业学业与职业规划.pptx"

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9 宽屏
prs.slide_height = Inches(7.5)

# ============================================================
# 辅助函数
# ============================================================

def add_blank_slide():
    """添加空白幻灯片"""
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def add_textbox(slide, left, top, width, height, text="", font_size=14,
                bold=False, color=DARK, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_rich_textbox(slide, left, top, width, height):
    """添加空文本框，返回 text_frame"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf

def add_paragraph(tf, text, font_size=14, bold=False, color=DARK, alignment=PP_ALIGN.LEFT,
                  space_after=Pt(6), space_before=Pt(0), level=0, font_name="Microsoft YaHei"):
    """在 text_frame 中添加段落"""
    if len(tf.paragraphs) == 1 and tf.paragraphs[0].text == "":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = space_after
    p.space_before = space_before
    p.level = level
    return p

def add_rect(slide, left, top, width, height, fill_color=RED, line_color=None):
    """添加矩形"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_section_title(slide, number, title, subtitle=""):
    """添加章节过渡页"""
    # 深灰背景
    add_rect(slide, 0, 0, 13.333, 7.5, DARK)
    # 左侧红色竖条
    add_rect(slide, 0.8, 2.2, 0.08, 3.0, RED)
    # 章节号
    add_textbox(slide, 1.3, 2.2, 2, 0.8, number, font_size=48, bold=True, color=RED)
    # 标题
    add_textbox(slide, 1.3, 3.0, 10, 1.2, title, font_size=36, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, 1.3, 4.2, 10, 0.8, subtitle, font_size=16, color=MID_GRAY)
    # 底部红线
    add_rect(slide, 1.3, 5.5, 3, 0.04, RED)

def add_slide_title(slide, title):
    """添加页面标题（顶部红色下划线）"""
    add_textbox(slide, 0.8, 0.35, 11.5, 0.7, title, font_size=28, bold=True, color=DARK)
    add_rect(slide, 0.8, 1.0, 2.5, 0.04, RED)

def add_footer(slide, page_num):
    """添加页脚"""
    add_rect(slide, 0, 7.2, 13.333, 0.3, LIGHT_BG)
    add_textbox(slide, 11.5, 7.2, 1.5, 0.3, str(page_num), font_size=9, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)
    add_textbox(slide, 0.5, 7.2, 4, 0.3, "西安建筑科技大学管理学院 · 工商管理专业", font_size=8, color=MID_GRAY)


# ============================================================
# 第1页：封面
# ============================================================
slide = add_blank_slide()
# 白色背景 + 底部红色条
add_rect(slide, 0, 6.8, 13.333, 0.08, RED)
add_rect(slide, 0, 6.88, 13.333, 0.62, DARK)
add_textbox(slide, 0.8, 7.0, 5, 0.3, "西安建筑科技大学管理学院 · 2026年5月 · 专题交流活动",
            font_size=9, color=MID_GRAY)
# 顶部左侧装饰
add_rect(slide, 0, 0, 0.08, 1.2, RED)
add_rect(slide, 0.08, 0, 13.25, 0.06, RED)
# 主标题
add_textbox(slide, 1.5, 2.2, 10, 1.4, "工商管理专业", font_size=48, bold=True, color=RED)
add_textbox(slide, 1.5, 3.5, 10, 1.2, "学业与职业规划", font_size=42, bold=True, color=DARK)
# 副标题
tf = add_rich_textbox(slide, 1.5, 4.8, 10, 0.6)
p = tf.paragraphs[0]
p.text = "—— 立足建筑行业 · 赋能管理未来"
p.font.size = Pt(18)
p.font.color.rgb = MID_GRAY
p.font.name = "Microsoft YaHei"
# 分隔线
add_rect(slide, 1.5, 5.6, 3, 0.04, RED)
# 小组信息
add_textbox(slide, 1.5, 5.9, 5, 0.5, "宿舍小组：追光者", font_size=14, bold=True, color=DARK)
# 隐形校徽位置提示
add_textbox(slide, 9.5, 0.4, 3, 0.4, "[ 校 徽 ]", font_size=11, color=MID_GRAY, alignment=PP_ALIGN.RIGHT)


# ============================================================
# 第2页：目录
# ============================================================
slide = add_blank_slide()
add_slide_title(slide, "目  录  CONTENTS")
add_rect(slide, 0.8, 1.2, 11.8, 0.5, LIGHT_BG)

sections = [
    ("01", "工商管理专业现状分析", "我们的优势 · 挑战 · 定位"),
    ("02", "专业未来发展研判", "数字化转型 · 绿色建筑 · 复合人才"),
    ("03", "大学生活与学习规划", "四年行动指南 · 技能提升路径"),
    ("04", "未来职业方向选择", "就业 · 升学 · 考公三条路径"),
    ("05", "小组分工与成员观点", "六种视角 · 一种使命"),
]

for i, (num, title, desc) in enumerate(sections):
    y = 1.9 + i * 1.0
    # 数字
    add_textbox(slide, 1.0, y, 1.2, 0.7, num, font_size=32, bold=True, color=RED)
    # 竖线
    add_rect(slide, 2.1, y + 0.08, 0.04, 0.55, RED)
    # 标题
    add_textbox(slide, 2.4, y + 0.05, 4, 0.4, title, font_size=18, bold=True, color=DARK)
    add_textbox(slide, 2.4, y + 0.4, 6, 0.3, desc, font_size=11, color=MID_GRAY)

add_footer(slide, 2)


# ============================================================
# 第3页：Part 1 过渡页
# ============================================================
slide = add_blank_slide()
add_section_title(slide, "01", "工商管理专业现状分析", "我们的优势 · 我们的挑战 · 我们的定位")
add_footer(slide, 3)


# ============================================================
# 第4页：学科优势
# ============================================================
slide = add_blank_slide()
add_slide_title(slide, "学科优势：建筑 + 管理的双重基因")
add_footer(slide, 4)

# 三列卡片
advantages = [
    ("🏛️ 行业底蕴深厚", [
        "专业源自1956年",
        "\"建筑经济组织与计划管理\"",
        "60+年历史传承",
        "国家级一流本科专业建设点",
        "工商管理一级学科硕士点支撑",
    ]),
    ("🔗 建筑·管理深度融合", [
        "五大特色微方向：",
        "城市建设与运营",
        "数据智能与商业决策",
        "供应链运营 · 市场营销 · HR",
        "含BIM/造价/国际工程课程",
    ]),
    ("🎯 行业影响力背书", [
        "近三年400+毕业生",
        "入职世界500强（中建/中铁等）",
        "管理学院就业率90%+",
        "工程管理2025届",
        "就业率高达 97.70%",
    ]),
]

for i, (title, items) in enumerate(advantages):
    x = 0.8 + i * 4.1
    # 卡片背景
    add_rect(slide, x, 1.3, 3.8, 5.4, LIGHT_BG)
    add_rect(slide, x, 1.3, 3.8, 0.06, RED)
    # 卡片标题
    add_textbox(slide, x + 0.3, 1.5, 3.2, 0.5, title, font_size=16, bold=True, color=RED)
    # 内容
    tf = add_rich_textbox(slide, x + 0.3, 2.2, 3.2, 4.2)
    for j, item in enumerate(items):
        add_paragraph(tf, item, font_size=12, bold=(j == 0), color=DARK if j > 0 else DARK,
                      space_after=Pt(4))
    add_footer(slide, 4)

# 底部引用
add_textbox(slide, 0.8, 6.75, 12, 0.35,
            "数据来源：西安建筑科技大学管理学院官网 (som.xauat.edu.cn) 2023-2025届就业统计",
            font_size=8, color=MID_GRAY)


# ============================================================
# 第5页：学科劣势
# ============================================================
slide = add_blank_slide()
add_slide_title(slide, "现实挑战：我们需要正视的问题")
add_footer(slide, 5)

# 表格布局：5行问题卡片
challenges = [
    ("课程广度 vs 深度", "管理学/经济学/财务/HR都学，但每门不够深入", "大三选择微方向聚焦，辅修BIM/造价课程"),
    ("建筑行业周期性", "房地产下行影响传统岗位，2025届工商管理就业率82.35%", "拓展数字化管理技能，不局限传统建筑企业"),
    ("技术壁垒不足", "相比工程管理/土木专业，缺乏硬核技术标签", "主动学习BIM基础、工程造价、Python数据分析"),
    ("实习竞争激烈", "优质建筑企业（陕建/中建）实习名额有限", "大二起参加商赛、积累项目经验，提升竞争力"),
    ("区域局限性", "西北地区建筑产业规模小于东部沿海", "善用校友网络（中西部建企认可度高）"),
]

# 表头
headers = ["挑战维度", "具体表现", "应对策略"]
for j, h in enumerate(headers):
    x_map = [0.8, 3.2, 7.5]
    w_map = [2.2, 4.1, 4.8]
    add_rect(slide, x_map[j], 1.3, w_map[j], 0.45, RED)
    add_textbox(slide, x_map[j] + 0.15, 1.32, w_map[j] - 0.3, 0.4, h,
                font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# 行
for i, (dim, problem, solution) in enumerate(challenges):
    y = 1.8 + i * 1.02
    bg = LIGHT_BG if i % 2 == 0 else WHITE
    x_map = [0.8, 3.2, 7.5]
    w_map = [2.2, 4.1, 4.8]
    texts = [dim, problem, solution]
    for j in range(3):
        add_rect(slide, x_map[j], y, w_map[j], 0.92, bg)
        add_textbox(slide, x_map[j] + 0.15, y + 0.1, w_map[j] - 0.3, 0.72, texts[j],
                    font_size=11, color=DARK)

# 底部
add_textbox(slide, 0.8, 6.8, 12, 0.3,
            "📌 短板不可怕，关键是用行动补齐 —— 提前2-3年规划执行",
            font_size=11, bold=True, color=RED)


# ============================================================
# 第6页：课程体系与行业适配度
# ============================================================
slide = add_blank_slide()
add_slide_title(slide, "课程地图：从课堂到岗位")
add_footer(slide, 6)

# 左侧：课程树状
tf = add_rich_textbox(slide, 0.8, 1.3, 5.5, 5.5)
courses_tree = [
    ("📐 通识教育", "高等数学 / 线性代数 / 概率论 / 土木工程概论 / 工程制图", 11, MID_GRAY),
    ("📊 学科基础", "管理学 / 微观经济学 / 运筹学 / 管理统计学 / 基础会计学", 11, BLUE),
    ("🎓 专业核心", "组织行为学 / 企业战略管理 / 财务管理 / 运营管理 / HRM / 市场营销", 11, DARK),
    ("🏗️ 行业特色", "技术经济学 / 房地产开发与经营 / 工程项目管理 / 建筑企业管理 / 国际工程投资策划", 11, RED),
    ("🤖 数智赋能", "大数据与商业智能决策 / 生成式AI商务应用 / 商务数据分析 / 数字化营销", 11, GREEN),
]
for label, content, size, color in courses_tree:
    add_paragraph(tf, f"{label}", font_size=13, bold=True, color=color, space_after=Pt(2))
    add_paragraph(tf, f"    {content}", font_size=10, color=MID_GRAY, space_after=Pt(10))

# 右侧：岗位映射表格
right_x = 6.8
add_textbox(slide, right_x, 1.3, 5.5, 0.4, "🏢 行业适配岗位映射", font_size=16, bold=True, color=DARK)

mapping = [
    ("项目管理+技术经济学", "建筑施工企业", "项目助理 → 项目经理"),
    ("房地产开发+市场营销", "房地产开发企业", "策划专员 → 投资拓展经理"),
    ("财务管理+经济法", "工程咨询企业", "成本控制 → 造价咨询"),
    ("人力资源管理+组织行为", "建筑企业总部", "管培生 → HRBP"),
    ("国际工程+跨文化管理", "海外工程企业", "国际商务经理"),
    ("数智管理+商务数据分析", "建筑科技企业", "数字化运营专员"),
]

for i, (course, industry, job) in enumerate(mapping):
    y = 1.85 + i * 0.85
    add_rect(slide, right_x, y, 5.8, 0.75, LIGHT_BG if i % 2 == 0 else WHITE)
    add_textbox(slide, right_x + 0.15, y + 0.02, 2.4, 0.35, course, font_size=10, bold=True, color=RED)
    add_textbox(slide, right_x + 0.15, y + 0.34, 2.4, 0.3, industry, font_size=9, color=BLUE)
    add_textbox(slide, right_x + 2.7, y + 0.1, 2.9, 0.5, job, font_size=10, bold=True, color=DARK)


# ============================================================
# 第7页：Part 2 过渡页
# ============================================================
slide = add_blank_slide()
add_section_title(slide, "02", "专业未来发展研判", "数字化转型 · 绿色建筑 · 复合型人才需求")
add_footer(slide, 7)


# ============================================================
# 第8页：三大变革
# ============================================================
slide = add_blank_slide()
add_slide_title(slide, "建筑行业管理岗位的三大变革")
add_footer(slide, 8)

changes = [
    ("🔵 数字化转型", BLUE, [
        "智慧工地市场 CAGR 超15%",
        "新岗位：数字孪生项目经理",
        "BIM总监 · IoT项目总监",
        "71%建筑业高管对前景乐观",
        "（毕马威2026全球调查）",
    ]),
    ("🟢 绿色低碳", GREEN, [
        "\"双碳\"政策推动绿色建筑",
        "占比目标提升至60%以上",
        "新岗位：碳排放管理师",
        "ESG合规经理",
        "绿建咨询市场规模达千亿级",
    ]),
    ("🟠 国际化拓展", ORANGE, [
        "\"一带一路\"沿线",
        "在建项目超3000个",
        "新岗位：国际工程商务经理",
        "FIDIC合同专员",
        "中国对外承包工程",
        "年营业额超1600亿美元",
    ]),
]

for i, (title, accent, items) in enumerate(changes):
    x = 0.8 + i * 4.1
    add_rect(slide, x, 1.3, 3.8, 5.0, LIGHT_BG)
    add_rect(slide, x, 1.3, 3.8, 0.06, accent)
    add_textbox(slide, x + 0.3, 1.5, 3.2, 0.5, title, font_size=16, bold=True, color=accent)
    tf = add_rich_textbox(slide, x + 0.3, 2.1, 3.2, 3.8)
    for item in items:
        bold = ("新岗位" in item) or ("CAGR" in item) or ("3000" in item) or ("千亿" in item) or ("1600" in item) or ("60%" in item)
        add_paragraph(tf, f"• {item}", font_size=12, bold=bold, color=DARK, space_after=Pt(5))

# 底部醒示
add_rect(slide, 0.8, 6.5, 11.8, 0.55, RED)
add_textbox(slide, 1.0, 6.5, 11.5, 0.55,
            "📊 复合型人才缺口：中间层\"懂技术+懂管理\"人才缺口高达 120万（2026行业调研）",
            font_size=14, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)


# ============================================================
# 第9页：人才画像
# ============================================================
slide = add_blank_slide()
add_slide_title(slide, "未来5年最具竞争力的管理人才画像")
add_footer(slide, 9)

# 四象限
# 左上卡片
add_rect(slide, 0.8, 1.3, 5.8, 2.4, LIGHT_BLUE)
add_textbox(slide, 1.1, 1.4, 5.3, 0.4, "★ 建筑科技产品经理（跨界稀缺型）", font_size=15, bold=True, color=BLUE)
tf = add_rich_textbox(slide, 1.1, 1.9, 5.3, 1.5)
for item in ["既懂建筑业务，又懂IT产品逻辑", "能定义智慧工地/BIM协同SaaS产品", "年薪区间：25-45万（一线城市）"]:
    add_paragraph(tf, f"• {item}", font_size=11, color=DARK, space_after=Pt(4))

# 右上卡片
add_rect(slide, 7.0, 1.3, 5.8, 2.4, LIGHT_GREEN)
add_textbox(slide, 7.3, 1.4, 5.3, 0.4, "★ 数字孪生项目经理（行业领军型）", font_size=15, bold=True, color=GREEN)
tf = add_rich_textbox(slide, 7.3, 1.9, 5.3, 1.5)
for item in ["能用BIM+GIS做全过程项目管理", "掌握EPC总承包+全过程工程咨询", "年薪区间：30-60万（海外项目更高）"]:
    add_paragraph(tf, f"• {item}", font_size=11, color=DARK, space_after=Pt(4))

# 左下卡片
add_rect(slide, 0.8, 4.0, 5.8, 2.4, LIGHT_ORANGE)
add_textbox(slide, 1.1, 4.1, 5.3, 0.4, "★ 工程商务经理（转型升级型）", font_size=15, bold=True, color=ORANGE)
tf = add_rich_textbox(slide, 1.1, 4.6, 5.3, 1.5)
for item in ["精通FIDIC合同+国际工程计价", "能用数据分析管控项目成本与风险", "年薪区间：18-35万（含驻外补贴）"]:
    add_paragraph(tf, f"• {item}", font_size=11, color=DARK, space_after=Pt(4))

# 右下卡片
add_rect(slide, 7.0, 4.0, 5.8, 2.4, LIGHT_RED)
add_textbox(slide, 7.3, 4.1, 5.3, 0.4, "⚠️ 传统资料员/文员（即将被替代）", font_size=15, bold=True, color=RED)
tf = add_rich_textbox(slide, 7.3, 4.6, 5.3, 1.5)
for item in ["AI+数字化流程已取代80%基础文职", "仅会抄写/归档将无竞争力", "→ 需升级为数字化档案管理岗"]:
    add_paragraph(tf, f"• {item}", font_size=11, color=DARK, space_after=Pt(4))

# 底部公式
add_rect(slide, 0.8, 6.6, 11.8, 0.45, DARK)
add_textbox(slide, 1.0, 6.6, 11.5, 0.45,
            "🔑 竞争力公式：建筑业务知识(40%) + 数据分析能力(30%) + 管理协调能力(30%)",
            font_size=13, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)


# ============================================================
# 第10页：SWOT
# ============================================================
slide = add_blank_slide()
add_slide_title(slide, "机遇 vs 挑战：我们的SWOT研判")
add_footer(slide, 10)

swot_data = [
    ("S 优势 Strengths", GREEN, [
        "✅ 国家级一流专业建设点",
        "✅ 60+年建筑经济管理底蕴",
        "✅ 校友网络在中西部建企认可度高",
        "✅ 数智化方向布局领先（AIGC课程）",
    ]),
    ("W 劣势 Weaknesses", ORANGE, [
        "❌ 课程广度大于深度",
        "❌ 技术硬技能需额外补充",
        "❌ 与工管/土木比专业辨识度略低",
        "❌ 区域（西北）建筑产业规模有限",
    ]),
    ("O 机会 Opportunities", BLUE, [
        "🌟 绿色建筑/双碳催生新型管理岗",
        "🌟 智慧工地/BIM推广急需复合人才",
        "🌟 \"一带一路\"国际工程管理需求",
        "🌟 全过程工程咨询政策推进",
    ]),
    ("T 威胁 Threats", RED, [
        "⚠️ 房地产下行影响传统岗位",
        "⚠️ AI替代部分基础管理职能",
        "⚠️ 东部高校毕业生竞争",
        "⚠️ 建筑行业周期性波动",
    ]),
]

for i, (title, accent, items) in enumerate(swot_data):
    col = i % 2
    row = i // 2
    x = 0.8 + col * 6.2
    y = 1.3 + row * 2.5
    add_rect(slide, x, y, 5.9, 2.2, WHITE)
    add_rect(slide, x, y, 5.9, 0.06, accent)
    add_textbox(slide, x + 0.2, y + 0.15, 5.5, 0.35, title, font_size=14, bold=True, color=accent)
    tf = add_rich_textbox(slide, x + 0.2, y + 0.55, 5.5, 1.5)
    for item in items:
        add_paragraph(tf, item, font_size=11, color=DARK, space_after=Pt(3))

# 策略
add_rect(slide, 0.8, 6.5, 11.8, 0.55, LIGHT_BG)
add_textbox(slide, 1.0, 6.5, 11.5, 0.55,
            "♟️ 策略：SO优先对接绿色建筑/国际工程实习  |  WO弥补技术短板（BIM+数据分析）  |  ST拓展建筑科技新赛道  |  WT考取行业证书增强硬实力",
            font_size=10, bold=True, color=DARK)


# ============================================================
# 第11页：Part 3 过渡页
# ============================================================
slide = add_blank_slide()
add_section_title(slide, "03", "大学生活与学习规划", "四年 · 四个阶段 · 一份成长路线图")
add_footer(slide, 11)


# ============================================================
# 第12页：四年规划时间轴
# ============================================================
slide = add_blank_slide()
add_slide_title(slide, "学业与能力提升四年路线图")
add_footer(slide, 12)

year_data = [
    ("📚 大一：夯实基础", GREEN, [
        "课程：通识课+数学三件套",
        "土木工程概论（了解行业）",
        "★ 考级：CET-4/6 四六级",
        "★ 计算机二级（MS Office）",
        "视野：阅读建筑行业概览书籍",
        "参加学院职业规划讲座",
    ]),
    ("📖 大二：专业筑基", BLUE, [
        "课程：管理学/经济学/会计/运筹",
        "组织行为学/战略管理",
        "★ BIM职业技能等级证书",
        "★ 初级造价员基础学习",
        "加入房地产/商赛社团",
        "寒假开始第一份短期见习",
    ]),
    ("💼 大三：方向聚焦", ORANGE, [
        "选定微方向深入学习",
        "国际工程/房地产/数智方向",
        "★ 暑期企业实习（陕建/中建/万科）",
        "★ 参加挑战杯/互联网+竞赛",
        "考研：数学+英语启动复习",
        "考公：行测+申论初步了解",
    ]),
    ("🎓 大四：冲刺收获", RED, [
        "毕业设计+毕业论文",
        "★ 秋招/春招正式签约",
        "★ 考研统考（12月）+复试（3月）",
        "★ 考公国考（11月）+省考（3月）",
        "PMP/二级建造师（满足条件）",
        "实习转正答辩",
    ]),
]

for i, (title, accent, items) in enumerate(year_data):
    x = 0.6 + i * 3.15
    add_rect(slide, x, 1.4, 2.95, 5.5, LIGHT_BG)
    add_rect(slide, x, 1.4, 2.95, 0.06, accent)
    add_textbox(slide, x + 0.15, 1.55, 2.7, 0.4, title, font_size=14, bold=True, color=accent)
    tf = add_rich_textbox(slide, x + 0.15, 2.1, 2.7, 4.5)
    for item in items:
        is_key = item.startswith("★")
        add_paragraph(tf, item, font_size=10, bold=is_key, color=RED if is_key else DARK,
                      space_after=Pt(4))

# 底部标语
add_rect(slide, 0.8, 7.05, 11.8, 0.35, DARK)
add_textbox(slide, 1.0, 7.05, 11.5, 0.35,
            "「大一不迷茫 → 大二有方向 → 大三能落地 → 大四见分晓」",
            font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)


# ============================================================
# 第13页：核心技能提升
# ============================================================
slide = add_blank_slide()
add_slide_title(slide, "五大核心能力培养路径")
add_footer(slide, 13)

skills_data = [
    ("🛠 项目管理软件", "Microsoft Project / Primavera P6", "能独立编制进度计划", "大三上学期"),
    ("📊 建筑行业数据分析", "Excel高级 / Python基础 + 管理统计学", "能做成本分析+数据看板", "大二下→大三"),
    ("📐 BIM基础认知", "Revit基础 + 广联达BIM5D", "能看懂模型+提取工程量", "大二暑假"),
    ("🤝 沟通协调与汇报", "课堂展示 + 竞赛路演 + 社团干部", "独立进行项目汇报", "全程培养"),
    ("📋 工程合同与招投标", "《建设工程合同管理》+标准招标文件", "能读懂合同核心条款", "大三上学期"),
]

for i, (skill, resource, target, timing) in enumerate(skills_data):
    y = 1.4 + i * 1.05
    add_rect(slide, 0.8, y, 11.8, 0.9, LIGHT_BG if i % 2 == 0 else WHITE)
    add_textbox(slide, 1.0, y + 0.05, 3, 0.35, skill, font_size=13, bold=True, color=DARK)
    add_textbox(slide, 1.0, y + 0.42, 4, 0.3, f"学习资源：{resource}", font_size=10, color=MID_GRAY)
    add_textbox(slide, 4.5, y + 0.2, 3, 0.35, f"目标：{target}", font_size=11, color=BLUE)
    add_textbox(slide, 8.5, y + 0.2, 3.5, 0.35, f"⏱ {timing}", font_size=11, bold=True, color=RED)

# 实践路径
add_textbox(slide, 0.8, 6.65, 3, 0.35, "实践积累路径：", font_size=12, bold=True, color=DARK)
path_text = (
    "竞赛：课堂学习 → 院赛选拔 → 校赛 → \"互联网+\"/\"挑战杯\"省赛/国赛    |    "
    "实习：大一参观 → 大二短期见习 → 大三暑期正式实习 → 大四转正    |    "
    "证书：CET4/6 → 计算机二级 → BIM初级 → 初级造价员/初级经济师"
)
add_textbox(slide, 0.8, 6.9, 11.8, 0.3, path_text, font_size=8, color=MID_GRAY)


# ============================================================
# 第14页：Part 4 过渡页
# ============================================================
slide = add_blank_slide()
add_section_title(slide, "04", "未来职业方向选择", "就业 · 升学 · 考公 —— 三条路径全面解析")
add_footer(slide, 14)


# ============================================================
# 第15页：路径一 —— 就业
# ============================================================
slide = add_blank_slide()
add_slide_title(slide, "路径一：建筑行业管理岗位就业")
add_footer(slide, 15)

jobs = [
    ("建筑施工", "管培生→项目助理→项目经理", "陕建集团 / 中建三局西北公司", "6-8K/月起"),
    ("房地产开发", "策划专员→投资拓展经理", "万科西北 / 碧桂园陕西", "7-9K/月起"),
    ("工程咨询", "招投标专员→成本控制", "中国建设咨询 / 西北工程设计院", "5-7K/月起"),
    ("建筑科技", "数字化运营 / BIM协调员", "广联达 / 品茗科技", "8-10K/月起"),
    ("国际工程", "商务助理→国际商务经理", "中建海外部 / 中国电建", "10-15K/月(驻外)"),
]

# 表头
job_headers = ["岗位类别", "晋升路径", "目标企业举例", "起薪参考（西北）"]
job_x = [0.8, 3.0, 6.2, 9.5]
job_w = [2.0, 3.0, 3.1, 2.8]
for j, (h, x, w) in enumerate(zip(job_headers, job_x, job_w)):
    add_rect(slide, x, 1.3, w, 0.45, RED)
    add_textbox(slide, x + 0.1, 1.32, w - 0.2, 0.4, h, font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

for i, (cat, path, company, salary) in enumerate(jobs):
    y = 1.8 + i * 0.7
    bg = LIGHT_BG if i % 2 == 0 else WHITE
    texts = [cat, path, company, salary]
    for j, (text, x, w) in enumerate(zip(texts, job_x, job_w)):
        add_rect(slide, x, y, w, 0.62, bg)
        c = RED if j == 0 else BLUE if j == 3 else DARK
        b = True if j == 0 else (True if j == 3 else False)
        add_textbox(slide, x + 0.1, y + 0.1, w - 0.2, 0.42, text,
                    font_size=11, bold=b, color=c, alignment=PP_ALIGN.CENTER if j == 3 else PP_ALIGN.LEFT)

# 就业准备时间线
add_textbox(slide, 0.8, 5.5, 5, 0.4, "就业准备时间线：", font_size=13, bold=True, color=DARK)
timeline = [
    ("大二下", "简历初稿\n了解行业企业"),
    ("大三上", "简历2.0\n企业开放日"),
    ("大三暑假", "暑期正式实习\n争取转正"),
    ("大四上", "秋招投递\n群面+专业面试"),
    ("大四下", "春招补录\n最终签约"),
]
for i, (time, desc) in enumerate(timeline):
    x = 0.8 + i * 2.4
    # 圆圈
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + 0.5), Inches(5.95), Inches(0.18), Inches(0.18))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RED
    shape.line.fill.background()
    # 连线
    if i < 4:
        add_rect(slide, x + 0.7, 6.02, 2.2, 0.03, RED)
    add_textbox(slide, x, 6.2, 2.2, 0.35, time, font_size=11, bold=True, color=RED)
    add_textbox(slide, x, 6.5, 2.2, 0.5, desc, font_size=9, color=DARK)

# 适合人群
add_rect(slide, 0.8, 7.1, 11.8, 0.3, LIGHT_BG)
add_textbox(slide, 1.0, 7.1, 11.5, 0.3,
            "👤 适合人群：愿意下项目/工地 · 抗压能力强 · 对建筑行业有热情 · 注重实践经验积累",
            font_size=10, color=DARK)


# ============================================================
# 第16页：路径二 —— 升学
# ============================================================
slide = add_blank_slide()
add_slide_title(slide, "路径二：升学深造 —— 理论进阶与学历提升")
add_footer(slide, 16)

# 三条路线
routes = [
    ("🏫 本校读研", GREEN, "管理科学与工程 / 工商管理(MBA)", "西安建筑科技大学", "数三 + 管理学原理", "省级重点学科 · 可提前联系导师 · 保研名额增加"),
    ("🎯 外校考研", BLUE, "建筑类强校管理专业", "同济大学 / 重庆大学 / 哈工大", "数三/数一 + 专业课", "备考强度较大 · 需大三启动 · 关注目标院校夏令营"),
    ("🌍 出国留学", ORANGE, "工程管理 / 房地产金融", "英国雷丁 / 香港理工 / 墨尔本大学", "IELTS 6.5+ + GPA 3.0+", "关注QS专业排名 · 准备语言考试 · 提前积累科研经历"),
]

for i, (title, accent, direction, schools, exam, extra) in enumerate(routes):
    y = 1.3 + i * 1.75
    add_rect(slide, 0.8, y, 11.8, 1.55, LIGHT_BG if i % 2 == 0 else WHITE)
    add_rect(slide, 0.8, y, 0.06, 1.55, accent)
    add_textbox(slide, 1.2, y + 0.1, 2.5, 0.35, title, font_size=16, bold=True, color=accent)
    add_textbox(slide, 3.8, y + 0.1, 4, 0.35, f"方向：{direction}", font_size=12, bold=True, color=DARK)
    add_textbox(slide, 3.8, y + 0.55, 4, 0.3, f"目标：{schools}", font_size=11, color=BLUE)
    add_textbox(slide, 8.5, y + 0.1, 3.5, 0.35, f"考试：{exam}", font_size=11, color=RED)
    add_textbox(slide, 1.2, y + 1.0, 10, 0.35, f"💡 {extra}", font_size=10, color=MID_GRAY)

# 考研时间线
add_textbox(slide, 0.8, 6.65, 5, 0.35, "考研备考时间线：", font_size=12, bold=True, color=DARK)
exam_timeline = [
    ("大二暑假", "确定考研方向\n了解目标院校"),
    ("大三上", "数学基础复习\n每日背单词50个"),
    ("大三寒假", "英语真题+词汇\n加入学习小组"),
    ("大三下", "专业课系统复习\n暑假留校集训"),
    ("大四上", "冲刺模考9-12月\n12月全国统考"),
]
for i, (time, desc) in enumerate(exam_timeline):
    x = 0.8 + i * 2.4
    add_textbox(slide, x, 6.95, 2.2, 0.2, time, font_size=10, bold=True, color=RED)
    add_textbox(slide, x, 7.15, 2.2, 0.3, desc, font_size=8, color=DARK)


# ============================================================
# 第17页：路径三 —— 考公/考编
# ============================================================
slide = add_blank_slide()
add_slide_title(slide, "路径三：公务员与事业单位 —— 公共服务赛道")
add_footer(slide, 17)

gov_jobs = [
    ("国考", "住建部 / 发改委", "高度对口", "较高 (1:100+)"),
    ("省考", "省住建厅 / 自然资源厅", "高度对口", "中等偏高"),
    ("市/区考", "市住建局 / 区建设局", "对口", "中等"),
    ("事业单位", "工程质量监督站 / 公共资源交易中心", "对口", "中等"),
    ("国企编制", "城投公司 / 保障房建设公司", "高度对口", "中等偏低"),
]

gov_headers = ["层级", "部门/岗位", "专业对口度", "竞争比参考"]
gov_x = [0.8, 2.8, 6.5, 9.2]
gov_w = [1.8, 3.5, 2.5, 3.0]
for j, (h, x, w) in enumerate(zip(gov_headers, gov_x, gov_w)):
    add_rect(slide, x, 1.3, w, 0.45, RED)
    add_textbox(slide, x + 0.1, 1.32, w - 0.2, 0.4, h, font_size=12, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

for i, (level, dept, fit, comp) in enumerate(gov_jobs):
    y = 1.8 + i * 0.65
    bg = LIGHT_BG if i % 2 == 0 else WHITE
    texts = [level, dept, fit, comp]
    for j, (text, x, w) in enumerate(zip(texts, gov_x, gov_w)):
        add_rect(slide, x, y, w, 0.55, bg)
        c = RED if j == 0 else GREEN if j == 2 else DARK
        add_textbox(slide, x + 0.1, y + 0.08, w - 0.2, 0.4, text,
                    font_size=11, bold=(j == 0 or j == 2), color=c)

# 优势
add_rect(slide, 0.8, 5.25, 11.8, 1.4, LIGHT_GREEN)
add_textbox(slide, 1.0, 5.3, 11.5, 0.35, "🏛️ 西建大考公核心优势", font_size=14, bold=True, color=GREEN)
advantages_gov = [
    "• 住建系统对西建大认可度较高（\"建筑老八校\"品牌效应）",
    "• 建筑类岗位限制专业，竞争比远低于\"三不限\"岗位",
    "• 校友网络广泛分布于西北地区住建系统",
]
tf = add_rich_textbox(slide, 1.0, 5.7, 11.5, 0.9)
for a in advantages_gov:
    add_paragraph(tf, a, font_size=11, color=DARK, space_after=Pt(3))

# 适合人群
add_textbox(slide, 0.8, 6.8, 11.8, 0.35,
            "👤 适合人群：追求稳定 · 对公共政策有兴趣 · 重视工作生活平衡 · 愿意长期在体制内发展",
            font_size=11, bold=True, color=DARK)


# ============================================================
# 第18页：小组分工与成员观点
# ============================================================
slide = add_blank_slide()
add_slide_title(slide, "小组分工与成员观点")
add_footer(slide, 18)

# 分工表
add_textbox(slide, 0.8, 1.3, 3, 0.35, "📋 分工表", font_size=14, bold=True, color=DARK)
roles = [
    ("成员A（组长）", "PPT框架+第一、二部分", "专业现状与趋势分析"),
    ("成员B", "第三部分：学习规划", "四年时间轴+技能提升路径"),
    ("成员C", "第四部分：职业方向", "三路径解析+数据搜集"),
    ("成员D", "设计美化+数据可视化", "PPT排版+图表制作"),
    ("成员E", "资料搜集+演讲备注", "课程数据+就业数据整理"),
    ("成员F", "审校+模拟演练", "内容审核+演讲演练"),
]

role_x = [0.8, 3.3, 7.0]
role_w = [2.3, 3.5, 4.8]
role_headers = ["成员", "分工模块", "核心贡献"]
for j, (h, x, w) in enumerate(zip(role_headers, role_x, role_w)):
    add_rect(slide, x, 1.75, w, 0.35, DARK)
    add_textbox(slide, x + 0.1, 1.75, w - 0.2, 0.35, h, font_size=10, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

for i, (member, module, contribution) in enumerate(roles):
    y = 2.15 + i * 0.35
    bg = LIGHT_BG if i % 2 == 0 else WHITE
    texts = [member, module, contribution]
    for j, (text, x, w) in enumerate(zip(texts, role_x, role_w)):
        add_rect(slide, x, y, w, 0.32, bg)
        c = RED if j == 0 else DARK
        add_textbox(slide, x + 0.1, y + 0.02, w - 0.2, 0.28, text,
                    font_size=9, bold=(j == 0), color=c)

# 成员观点
add_textbox(slide, 0.8, 4.4, 4, 0.35, "💬 成员个人观点摘录", font_size=14, bold=True, color=DARK)

opinions = [
    ("🟡 务实就业派 (成员A)", "\"我选择工商管理就是想进建筑行业做管理。陕建、中建对我们学校认可度很高，打算大三去项目上实习，从管培生做起，目标5-8年做到项目经理。\""),
    ("🟢 考研深造派 (成员B)", "\"本科工商管理竞争力不够突出，计划大三开始备考，目标同济大学管理科学与工程，读研期间专攻工程项目管理方向，毕业选择面会宽很多。\""),
    ("🟠 考公稳定派 (成员C)", "\"我对公共政策方向更感兴趣，打算参加国考或省考，目标住建局或发改委。\'建筑老八校\'的背景在专业限制岗位上有明显优势。\""),
    ("🔵 跨界转型派 (成员D)", "\"想尝试建筑科技方向。传统建筑业在转型数字化，我准备自学Python，目标是进广联达这类公司做产品运营或商务数据分析。\""),
    ("🟣 国际工程派 (成员E)", "\"英语底子不错，对国际工程管理感兴趣。学校开有\'一带一路市场行情\'和\'跨文化管理\'课程，想去中建海外部门做国际商务岗。\""),
    ("🔴 探索观望派 (成员F)", "\"还没完全确定方向，打算大二多参加商赛和社团活动，大三实习后再做决定。先保证不挂科+考过四六级，边走边看，不给自己设限。\""),
]

for i, (label, quote) in enumerate(opinions):
    col = i % 3
    row = i // 3
    x = 0.8 + col * 4.1
    y = 4.85 + row * 1.15
    add_rect(slide, x, y, 3.9, 1.0, LIGHT_BG)
    add_rect(slide, x, y, 3.9, 0.04, RED)
    add_textbox(slide, x + 0.15, y + 0.08, 3.6, 0.25, label, font_size=10, bold=True, color=DARK)
    add_textbox(slide, x + 0.15, y + 0.32, 3.6, 0.62, quote, font_size=8.5, color=DARK)


# ============================================================
# 第19页：结语
# ============================================================
slide = add_blank_slide()
# 深色背景
add_rect(slide, 0, 0, 13.333, 7.5, DARK)
add_rect(slide, 0.8, 1.5, 0.08, 4.5, RED)

# 主标题
add_textbox(slide, 1.3, 1.5, 10, 1.0, "规划未来 · 行则将至", font_size=40, bold=True, color=WHITE)

# 结语内容
conclusion_lines = [
    "今天的规划，是明天的底气。",
    "",
    "工商管理是一个\"上限很高\"的专业——它不会限定你的职业边界，",
    "但也不会替你铺好每一步路。",
    "",
    "在西安建筑科技大学管理学院，",
    "我们拥有 60年建筑经济管理的历史积淀，",
    "拥有 国家级一流专业 的平台支撑，",
    "拥有 \"建筑+管理+数智\" 的独特定位。",
    "",
    "无论选择就业、升学还是考公，",
    "提前2-3年规划，一步一个脚印执行，",
    "就是我们最大的竞争力。",
]
tf = add_rich_textbox(slide, 1.3, 2.7, 10.5, 3.5)
for line in conclusion_lines:
    if line == "":
        add_paragraph(tf, "", font_size=4, color=WHITE, space_after=Pt(0))
    else:
        add_paragraph(tf, line, font_size=15, color=MID_GRAY if "上限" in line or "不会" in line else WHITE,
                      space_after=Pt(4))

# 底部
add_rect(slide, 1.3, 6.3, 4, 0.04, RED)
add_textbox(slide, 1.3, 6.5, 10, 0.5, "感谢聆听！", font_size=18, bold=True, color=WHITE)
add_textbox(slide, 1.3, 6.85, 10, 0.35,
            "西安建筑科技大学管理学院 · 工商管理专业 · 宿舍小组全体成员 · 2026年5月",
            font_size=10, color=MID_GRAY)

# ============================================================
# 保存文件
# ============================================================
os.makedirs(os.path.dirname(OUTPUT_PATH), exist_ok=True)
prs.save(OUTPUT_PATH)
print(f"PPT generated: {OUTPUT_PATH}")
print(f"Total slides: {len(prs.slides)}")
