"""
工商管理专业学业与职业规划 PPT V2 —— 视觉优化版
python-pptx 完整生成
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
import os

# ============================================================
# DESIGN TOKENS
# ============================================================
RED    = RGBColor(0xC4, 0x1E, 0x2A)
DARK   = RGBColor(0x2D, 0x2D, 0x2D)
GRAY   = RGBColor(0x70, 0x70, 0x70)
LIGHT  = RGBColor(0xF2, 0xF2, 0xF2)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
BLUE   = RGBColor(0x1A, 0x73, 0xE8)
GREEN  = RGBColor(0x13, 0x7B, 0x33)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
TEAL   = RGBColor(0x00, 0x96, 0x88)
PURPLE = RGBColor(0x7B, 0x1F, 0xA2)
CREAM  = RGBColor(0xFD, 0xF5, 0xEE)

OUT = r"D:\develop\cursor Pro\shao-xiaoli-attendance\工商管理专业学业与职业规划.pptx"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── helpers ─────────────────────────────────────────────
def _box(slide, l, t, w, h, fill=None, line=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l),Inches(t),Inches(w),Inches(h))
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line; s.line.width = Pt(1)
    else:
        s.line.fill.background()
    return s

def _round_box(slide, l, t, w, h, fill=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(l),Inches(t),Inches(w),Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill if fill else WHITE
    s.line.fill.background()
    return s

def _tb(slide, l, t, w, h, txt="", fs=16, bold=False, color=DARK, align=PP_ALIGN.LEFT, font="Microsoft YaHei"):
    tb = slide.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = txt; p.font.size = Pt(fs); p.font.bold = bold
    p.font.color.rgb = color; p.font.name = font; p.alignment = align
    return tf

def _rbox(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(Inches(l),Inches(t),Inches(w),Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    return tf

def _p(tf, txt, fs=14, bold=False, color=DARK, align=PP_ALIGN.LEFT, sa=Pt(6), sb=Pt(0), font="Microsoft YaHei"):
    if len(tf.paragraphs)==1 and tf.paragraphs[0].text=="":
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = txt; p.font.size = Pt(fs); p.font.bold = bold
    p.font.color.rgb = color; p.font.name = font; p.alignment = align
    p.space_after = sa; p.space_before = sb
    return p

def _circle(slide, l, t, d, fill=RED):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l),Inches(t),Inches(d),Inches(d))
    s.fill.solid(); s.fill.fore_color.rgb = fill; s.line.fill.background()
    return s

def _section(slide, num, title, sub=""):
    _box(slide, 0,0,13.333,7.5, DARK)
    _box(slide, 1.2, 2.5, 0.1, 3.0, RED)
    _tb(slide, 1.8, 2.5, 2, 1.0, num, fs=56, bold=True, color=RED)
    _tb(slide, 1.8, 3.4, 10, 1.0, title, fs=38, bold=True, color=WHITE)
    if sub: _tb(slide, 1.8, 4.4, 10, 0.6, sub, fs=16, color=GRAY)

def _page_head(slide, title):
    _tb(slide, 1.0, 0.45, 11, 0.7, title, fs=28, bold=True, color=DARK)
    _box(slide, 1.0, 1.15, 3.0, 0.06, RED)

def _footer(slide, n):
    _box(slide, 0, 7.2, 13.333, 0.3, LIGHT)
    _tb(slide, 1.0, 7.22, 6, 0.25, "西安建筑科技大学管理学院 · 工商管理专业 · 2026.05", fs=8, color=GRAY)
    _tb(slide, 11.5, 7.22, 1.5, 0.25, str(n), fs=9, color=RED, bold=True, align=PP_ALIGN.RIGHT)

def _card(slide, l, t, w, h, accent, title, lines, title_fs=15, body_fs=12):
    _round_box(slide, l, t, w, h, WHITE)
    _box(slide, l, t, w, 0.07, accent)
    _tb(slide, l+0.35, t+0.25, w-0.7, 0.45, title, fs=title_fs, bold=True, color=accent)
    tf = _rbox(slide, l+0.35, t+0.8, w-0.7, h-1.1)
    for line in lines:
        _p(tf, line, fs=body_fs, sa=Pt(7))
    return tf

# ============================================================
# SLIDE 1 — COVER
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
_box(sl, 0, 0, 13.333, 7.5, WHITE)
# top accent bar
_box(sl, 0, 0, 13.333, 0.1, RED)
# left big red block
_box(sl, 0, 2.5, 0.12, 3.5, RED)
# bottom bar
_box(sl, 0, 7.15, 13.333, 0.35, DARK)
_tb(sl, 1.2, 7.17, 10, 0.3, "西安建筑科技大学管理学院   |   专题交流活动   |   2026年5月9日", fs=10, color=GRAY)
# title
_tb(sl, 1.6, 2.7, 10, 1.2, "工商管理专业", fs=56, bold=True, color=RED)
_tb(sl, 1.6, 3.8, 10, 1.1, "学业与职业规划", fs=48, bold=True, color=DARK)
# subtitle
_tb(sl, 1.6, 5.0, 8, 0.5, "——  立足建筑行业  ·  赋能管理未来", fs=18, color=GRAY)
_box(sl, 1.6, 5.65, 4, 0.05, RED)
_tb(sl, 1.6, 5.95, 5, 0.45, "宿 舍 小 组 ：追 光 者", fs=16, bold=True, color=DARK)


# ============================================================
# SLIDE 2 — AGENDA
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
_page_head(sl, "目  录  CONTENTS")

items = [
    ("01", "工商管理专业现状分析", "优势 · 挑战 · 定位"),
    ("02", "专业未来发展研判", "数字化 · 绿色建筑 · 复合人才"),
    ("03", "大学生活与学习规划", "四年路线图 · 技能路径"),
    ("04", "未来职业方向选择", "就业 / 升学 / 考公"),
    ("05", "小组分工与成员观点", "六种视角 · 一种使命"),
]
for i,(num,ti,desc) in enumerate(items):
    y = 1.6 + i*1.1
    _box(sl, 2.5, y+0.15, 8.5, 0.95, LIGHT if i%2==0 else None)
    _tb(sl, 2.8, y+0.15, 1.5, 0.75, num, fs=36, bold=True, color=RED)
    _box(sl, 4.2, y+0.25, 0.06, 0.55, RED)
    _tb(sl, 4.6, y+0.15, 4.5, 0.45, ti, fs=20, bold=True, color=DARK)
    _tb(sl, 4.6, y+0.54, 5, 0.3, desc, fs=12, color=GRAY)
_footer(sl, 2)


# ============================================================
# SLIDE 3 — SECTION 1
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
_section(sl, "01", "工商管理专业现状分析", "我们的优势 · 我们的挑战 · 我们的定位")


# ============================================================
# SLIDE 4 — ADVANTAGES (3 cards + big numbers)
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
_page_head(sl, "学科优势：建筑 + 管理的双重基因")

cards = [
    (RED, "60+", "年历史传承", [
        "专业源自1956年\"建筑经济组织与计划管理\"",
        "国家级一流本科专业建设点",
        "工商管理一级学科硕士点支撑",
        "陕西省专业综合改革试点",
    ]),
    (TEAL, "5", "大特色微方向", [
        "城市建设与运营（房开+项目管理）",
        "数据智能与商业决策（AIGC+BI）",
        "市场营销 / 人力资源 / 供应链运营",
        "含国际工程投资策划等前沿课程",
    ]),
    (BLUE, "400+", "人入职世界500强", [
        "近三年入职中建、中铁等500强企业",
        "管理学院整体就业率稳定90%以上",
        "工程管理专业2025届就业率97.70%",
        "校友网络遍布中西部地区建企",
    ]),
]

for i,(accent,big_num, tag, lines) in enumerate(cards):
    x = 1.0 + i*4.0
    _round_box(sl, x, 1.5, 3.7, 5.2, WHITE)
    _box(sl, x, 1.5, 3.7, 0.08, accent)
    # big number
    _tb(sl, x+0.3, 1.8, 1.5, 0.7, big_num, fs=40, bold=True, color=accent)
    _tb(sl, x+1.7, 2.05, 1.8, 0.45, tag, fs=15, bold=True, color=DARK)
    _box(sl, x+0.3, 2.55, 3.1, 0.03, accent)
    tf = _rbox(sl, x+0.3, 2.8, 3.2, 3.5)
    for li in lines:
        _p(tf, f"✦ {li}", fs=13, sa=Pt(12))
_footer(sl, 4)


# ============================================================
# SLIDE 5 — CHALLENGES (table redesigned, bigger)
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
_page_head(sl, "现实挑战：我们需要正视的问题")

# 5 challenge rows - simpler layout, more impactful
ch = [
    ("课程广度多于深度", "❖ 管理学/经济/财务/HR都学，但每门不精   →   大三选择微方向聚焦，辅修BIM/造价"),
    ("建筑行业周期性波动", "❖ 房地产下行影响传统岗位，工商管理就业率82.35%   →   拓展数字化技能，不只盯着传统建企"),
    ("技术硬壁垒不足", "❖ 相比工管/土木缺乏硬核技术标签   →   主动学BIM基础、Python数据分析、工程造价"),
    ("优质实习竞争激烈", "❖ 陕建/中建等企业实习名额有限   →   大二起参加商赛+项目经验，提升简历含金量"),
    ("西北区域产业规模有限", "❖ 建筑体量小于东部沿海   →   善用校友网络+瞄准中西部建企总部/区域公司"),
]
for i,(title, detail) in enumerate(ch):
    y = 1.5 + i*1.08
    _round_box(sl, 1.0, y, 11.3, 0.92, WHITE)
    _circle(sl, 1.25, y+0.22, 0.3, RED)
    _tb(sl, 1.75, y+0.08, 4, 0.4, title, fs=16, bold=True, color=DARK)
    _tb(sl, 1.75, y+0.48, 10, 0.35, detail, fs=12, color=DARK)

_tb(sl, 1.0, 6.9, 11, 0.3, "▸ 短板不可怕，关键是用行动补齐 —— 提前2-3年规划执行", fs=14, bold=True, color=RED)
_footer(sl, 5)


# ============================================================
# SLIDE 6 — COURSE TO JOB MAP (split layout)
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
_page_head(sl, "课程地图：从课堂到岗位")

# Left - course tree
_box(sl, 1.0, 1.5, 5.5, 5.3, LIGHT)
_tb(sl, 1.3, 1.6, 4, 0.4, "📚 四层课程体系", fs=16, bold=True, color=RED)

tree = [
    ("通识基础", "高数 / 线代 / 概率论 / 工图 / 土木概论", GRAY),
    ("学科基础", "管理学 / 微观经济学 / 运筹学 / 管理统计 / 会计学", BLUE),
    ("专业核心", "组织行为学 / 战略管理 / 财务管理 / 运营 / HRM / 市场营销", DARK),
    ("行业特色", "技术经济学 / 房开经营 / 工程项目管理 / 建筑企业管理", RED),
    ("数智前沿", "大数据商业决策 / 生成式AI / 商务数据分析 / 数字化营销", GREEN),
]
tf = _rbox(sl, 1.3, 2.2, 5.0, 4.2)
for label, content, color in tree:
    _p(tf, f"▌{label}", fs=14, bold=True, color=color, sa=Pt(3))
    _p(tf, f"   {content}", fs=10, color=GRAY, sa=Pt(14))

# Right - job mapping
_box(sl, 6.9, 1.5, 5.5, 5.3, WHITE)
_tb(sl, 7.2, 1.6, 4, 0.4, "🏢 课岗精准对应", fs=16, bold=True, color=BLUE)

mapping = [
    ("项目+技术经济学", "施工企业 → 项目经理"),
    ("房开经营+营销", "房企 → 投资拓展/策划"),
    ("财务管理+经济法", "工程咨询 → 成本控制"),
    ("HRM+组织行为学", "建企总部 → HRBP"),
    ("国际工程+跨文化", "海外工程 → 商务经理"),
    ("数智管理+数据分析", "建筑科技 → 数字运营"),
]
for i,(course,job) in enumerate(mapping):
    y = 2.2 + i*0.75
    _circle(sl, 7.2, y+0.1, 0.2, RED)
    _tb(sl, 7.6, y-0.02, 3, 0.35, course, fs=12, bold=True, color=RED)
    _tb(sl, 7.6, y+0.3, 4.5, 0.3, job, fs=13, color=DARK)
_footer(sl, 6)


# ============================================================
# SLIDE 7 — SECTION 2
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
_section(sl, "02", "专业未来发展研判", "数字化转型 · 绿色建筑 · 复合型人才需求")


# ============================================================
# SLIDE 8 — THREE TRENDS (big visual cards)
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
_page_head(sl, "建筑行业管理岗位的三大变革")

trends = [
    (BLUE, "数字化转型", [
        "智慧工地市场规模年增超15%",
        "数字孪生项目经理 / BIM总监",
        "IoT项目总监等新岗位涌现",
        "71%建筑业高管对前景乐观",
        "—— 毕马威2026全球建筑业调查",
    ]),
    (GREEN, "绿色低碳", [
        "\"双碳\"政策推动绿色建筑占比",
        "到2030年目标提升至60%以上",
        "碳排放管理师 / ESG合规经理",
        "绿色建筑咨询市场规模千亿级",
        "—— 国家发改委绿色建筑行动方案",
    ]),
    (ORANGE, "国际化拓展", [
        "\"一带一路\"沿线在建项目超3000个",
        "国际工程商务经理 / FIDIC合同专员",
        "对外承包工程年营业额超1600亿美元",
        "全过程工程咨询政策全面推进",
        "—— 中国对外承包工程商会2026数据",
    ]),
]
for i,(accent,title,lines) in enumerate(trends):
    x = 1.0 + i*4.0
    _round_box(sl, x, 1.5, 3.7, 4.6, WHITE)
    _box(sl, x, 1.5, 3.7, 1.1, accent)
    _tb(sl, x+0.3, 1.6, 3.1, 0.9, title, fs=22, bold=True, color=WHITE)
    tf = _rbox(sl, x+0.3, 2.9, 3.1, 3.0)
    for li in lines:
        _p(tf, f"● {li}", fs=12, sa=Pt(8), color=DARK if not li.startswith("——") else GRAY)

# gap alert bar
_box(sl, 1.0, 6.35, 11.3, 0.55, RED)
_tb(sl, 1.3, 6.35, 10.7, 0.55,
    "人才缺口：中间层\"懂技术+懂管理\"的复合型人才缺口高达 120万 —— 不是人多，是符合新要求的人太少",
    fs=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
_footer(sl, 8)


# ============================================================
# SLIDE 9 — TALENT PORTRAIT (4 quadrants)
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
_page_head(sl, "未来 5 年最具竞争力的管理人才画像")

quads = [
    (1.0, 1.5, BLUE, "建筑科技产品经理", "跨界稀缺型", "懂建筑业务+IT产品逻辑\n定义智慧工地/BIM SaaS产品\n年薪：25-45万（一线城市）"),
    (7.0, 1.5, GREEN, "数字孪生项目经理", "行业领军型", "用BIM+GIS做全过程项目管理\nEPC总承包+全过程工程咨询\n年薪：30-60万（海外更高）"),
    (1.0, 4.0, ORANGE, "工程商务经理", "转型升级型", "精通FIDIC合同+国际计价\n用数据分析管控成本与风险\n年薪：18-35万（含驻外补贴）"),
    (7.0, 4.0, RED, "数字化档案/资料管理", "适应升级型", "AI+数字化流程已替代80%文职\n需掌握电子签章/一键组卷/合规\n传统\"抄抄写写\"将被淘汰"),
]
for (x,y,accent,title,tag,desc) in quads:
    _round_box(sl, x, y, 5.6, 2.2, WHITE)
    _box(sl, x, y, 0.1, 2.2, accent)
    _tb(sl, x+0.35, y+0.15, 3, 0.35, title, fs=17, bold=True, color=DARK)
    _box(sl, x+0.35, y+0.52, 2.5, 0.35, accent)
    _tb(sl, x+0.5, y+0.53, 2.3, 0.32, tag, fs=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _tb(sl, x+0.35, y+1.05, 5, 1.0, desc, fs=11, color=DARK)

# formula
_box(sl, 1.0, 6.45, 11.3, 0.48, DARK)
_tb(sl, 1.3, 6.45, 10.7, 0.48,
    "▸ 竞争力公式：建筑业务知识 (40%)  ＋  数据分析能力 (30%)  ＋  管理协调能力 (30%)",
    fs=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
_footer(sl, 9)


# ============================================================
# SLIDE 10 — SWOT (bigger, cleaner)
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
_page_head(sl, "机遇 vs 挑战：SWOT 研判")

swot = [
    (1.0, 1.5, GREEN,  "S  优势  Strengths",
     ["✦ 国家级一流专业建设点 | 60年建筑经济管理底蕴",
      "✦ 校友网络在中西部建企认可度高 | 数智化方向（AIGC）布局领先"]),
    (7.0, 1.5, ORANGE, "W  劣势  Weaknesses",
     ["✦ 课程广度大于深度 | 技术硬技能需额外补充",
      "✦ 与工管/土木比专业辨识度略低 | 西北区域产业规模有限"]),
    (1.0, 3.6, BLUE,   "O  机会  Opportunities",
     ["✦ 绿色建筑/双碳催生新型管理岗 | BIM推广急需复合人才",
      "✦ \"一带一路\"国际工程管理需求 | 全过程工程咨询政策推进"]),
    (7.0, 3.6, RED,    "T  威胁  Threats",
     ["✦ 房地产下行影响传统岗位 | AI替代部分基础管理职能",
      "✦ 东部高校毕业生竞争激烈 | 建筑行业周期性波动"]),
]
for (x,y,accent,title,lines) in swot:
    _round_box(sl, x, y, 5.6, 1.75, WHITE)
    _box(sl, x, y, 0.08, 1.75, accent)
    _tb(sl, x+0.3, y+0.12, 5, 0.4, title, fs=18, bold=True, color=accent)
    tf = _rbox(sl, x+0.3, y+0.65, 5, 0.9)
    for li in lines:
        _p(tf, li, fs=12, sa=Pt(6))

# strategy
_strats = [
    ("SO", "优先对接绿建/国际工程实习", GREEN),
    ("WO", "弥补技术短板（BIM+数据）", ORANGE),
    ("ST", "拓展建筑科技新赛道", BLUE),
    ("WT", "考取行业证书增强硬实力", RED),
]
for i,(tag,desc,color) in enumerate(_strats):
    x = 1.0 + i*3.0
    _round_box(sl, x, 5.65, 2.7, 0.6, color)
    _tb(sl, x+0.15, 5.68, 2.4, 0.54, f"{tag}策略：{desc}", fs=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

_footer(sl, 10)


# ============================================================
# SLIDE 11 — SECTION 3
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
_section(sl, "03", "大学生活与学习规划", "四年 · 四个阶段 · 一份成长路线图")


# ============================================================
# SLIDE 12 — 4 YEAR TIMELINE (horizontal, bigger)
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
_page_head(sl, "学业与能力提升四年路线图")

years = [
    (GREEN,  "大 一", "夯实基础", [
        "通识课 + 数学三件套",
        "土木工程概论（了解行业）",
        "⏺ CET-4/6 英语四六级",
        "⏺ 计算机二级（MS Office）",
        "阅读建筑行业概览书籍",
        "参加学院职业规划讲座",
    ]),
    (BLUE,   "大 二", "专业筑基", [
        "管理/经济/会计/运筹学",
        "组织行为学 / 战略管理",
        "⏺ BIM职业技能等级证书",
        "⏺ 初级造价员基础学习",
        "加入商赛/房地产社团",
        "寒假短期企业见习起步",
    ]),
    (ORANGE, "大 三", "方向聚焦", [
        "选定微方向深入学习",
        "国际工程/房地产/数智方向",
        "⏺ 暑期实习（陕建/中建/万科）",
        "⏺ 挑战杯/互联网+竞赛",
        "考研数学+英语基础启动",
        "考公行测+申论初步了解",
    ]),
    (RED,    "大 四", "冲刺收获", [
        "毕业设计 + 毕业论文",
        "⏺ 秋招签约 / 春招补录",
        "⏺ 考研统考12月 + 复试3月",
        "⏺ 考公国考11月 + 省考3月",
        "PMP/二级建造师（达标报考）",
        "实习转正答辩 / 三方落定",
    ]),
]
# header row
for i,(accent,label,sub,_) in enumerate(years):
    x = 0.8 + i*3.1
    _box(sl, x, 1.45, 2.95, 1.1, accent)
    _tb(sl, x+0.2, 1.5, 2.6, 0.55, label, fs=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    _tb(sl, x+0.2, 2.05, 2.6, 0.35, sub, fs=13, color=WHITE, align=PP_ALIGN.CENTER)

for i,(accent,label,sub,items) in enumerate(years):
    x = 0.8 + i*3.1
    _round_box(sl, x, 2.7, 2.95, 3.85, LIGHT)
    tf = _rbox(sl, x+0.2, 2.85, 2.6, 3.5)
    for it in items:
        key = it.startswith("⏺")
        _p(tf, it, fs=11.5, bold=key, color=RED if key else DARK, sa=Pt(8))

_box(sl, 1.0, 6.75, 11.3, 0.35, DARK)
_tb(sl, 1.2, 6.75, 11, 0.35,
    "「大一不迷茫 → 大二有方向 → 大三能落地 → 大四见分晓」",
    fs=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
_footer(sl, 12)


# ============================================================
# SLIDE 13 — SKILLS PATH
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
_page_head(sl, "五大核心能力培养路径")

skills = [
    (RED,    "项 目 管 理 软 件", "MS Project / Primavera P6", "能独立编制进度计划", "大三上学期"),
    (BLUE,   "建 筑 行 业 数 据 分 析", "Excel高级 + Python基础", "能做成本分析+数据看板", "大二下 → 大三"),
    (TEAL,   "BIM 基 础 认 知", "Revit基础 + 广联达BIM5D", "能看懂模型+提取工程量", "大二暑假"),
    (GREEN,  "沟 通 协 调 与 汇 报", "课堂展示 + 竞赛路演 + 社团", "独立进行项目汇报", "全程培养"),
    (ORANGE, "工 程 合 同 与 招 投 标", "建设工程合同管理 + 标准招标文件", "能读懂合同核心条款", "大三上学期"),
]
for i,(accent,skill,res,target,time) in enumerate(skills):
    y = 1.55 + i*1.02
    _round_box(sl, 1.0, y, 11.3, 0.88, WHITE)
    _box(sl, 1.0, y, 0.1, 0.88, accent)
    _tb(sl, 1.35, y+0.08, 3, 0.38, skill, fs=15, bold=True, color=DARK)
    _tb(sl, 1.35, y+0.48, 5, 0.3, f"📖 {res}", fs=11, color=GRAY)
    _tb(sl, 5.5, y+0.15, 4, 0.35, f"🎯 目标：{target}", fs=13, bold=True, color=BLUE)
    _tb(sl, 10.0, y+0.15, 3, 0.35, f"⏱ {time}", fs=13, bold=True, color=RED)

_tb(sl, 1.0, 6.75, 11, 0.35,
    "▸ 竞赛进阶：课堂学习 → 院赛选拔 → 校赛 → \"互联网+\"/\"挑战杯\"省赛/国赛   |   实习进阶：大一参观 → 大二见习 → 大三暑期实习 → 大四转正",
    fs=11, color=DARK)
_footer(sl, 13)


# ============================================================
# SLIDE 14 — SECTION 4
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
_section(sl, "04", "未来职业方向选择", "就业 · 升学 · 考公 —— 三条路径全面解析")


# ============================================================
# SLIDE 15 — PATH 1: JOBS
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
_page_head(sl, "路径一：建筑行业管理岗位就业")

jobs = [
    (RED,   "建筑施工", "管培生 → 项目助理 → 项目经理", "陕建集团 / 中建三局", "6-8K/月"),
    (BLUE,  "房地产开发", "策划专员 → 投资拓展经理", "万科西北 / 碧桂园陕西", "7-9K/月"),
    (TEAL,  "工程咨询", "招投标专员 → 成本控制", "中国建设咨询 / 西北设计院", "5-7K/月"),
    (GREEN, "建筑科技", "数字化运营 / BIM协调员", "广联达 / 品茗科技", "8-10K/月"),
    (ORANGE,"国际工程", "商务助理 → 国际商务经理", "中建海外部 / 中国电建", "10-15K/月"),
]

for i,(accent,cat,path,org,sal) in enumerate(jobs):
    y = 1.5 + i*0.85
    _round_box(sl, 1.0, y, 11.3, 0.7, WHITE)
    _circle(sl, 1.25, y+0.2, 0.3, accent)
    _tb(sl, 1.75, y+0.06, 2.8, 0.35, cat, fs=15, bold=True, color=accent)
    _tb(sl, 4.5, y+0.06, 3.5, 0.3, path, fs=12, color=DARK)
    _tb(sl, 4.5, y+0.36, 3.5, 0.25, org, fs=10, color=GRAY)
    _tb(sl, 9.8, y+0.1, 2.3, 0.4, sal, fs=15, bold=True, color=accent, align=PP_ALIGN.CENTER)

_tb(sl, 1.0, 5.9, 4, 0.35, "▸ 就业准备时间线", fs=15, bold=True, color=DARK)
tl = ["大二下\n简历初稿", "大三上\n企业开放日", "大三暑假\n正式实习", "大四上\n秋招投递", "大四下\n最终签约"]
for i,item in enumerate(tl):
    x = 1.0 + i*2.4
    _circle(sl, x+0.7, 6.35, 0.25, RED)
    if i<4: _box(sl, x+1.0, 6.45, 2.2, 0.05, RED)
    _tb(sl, x, 6.7, 2.2, 0.4, item, fs=10, bold=True, color=DARK, align=PP_ALIGN.CENTER)

_footer(sl, 15)


# ============================================================
# SLIDE 16 — PATH 2: GRAD SCHOOL
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
_page_head(sl, "路径二：升学深造 —— 理论进阶与学历提升")

routes = [
    (GREEN,  "本校读研", "管理科学与工程 / 工商管理(MBA)", "西安建筑科技大学", "数三 + 管理学原理", "省级重点学科 | 可提前联系导师 | 关注保研政策"),
    (BLUE,   "外校考研", "建筑类强校管理专业", "同济大学 / 重庆大学 / 哈工大", "数三/数一 + 专业课", "大三启动备考 | 关注目标院校夏令营 | 专业课针对性准备"),
    (ORANGE, "出国留学", "工程管理 / 房地产金融", "英国雷丁 / 香港理工 / 墨尔本大学", "IELTS 6.5+ + GPA 3.0+", "关注QS专业排名 | 提前准备语言 | 积累科研经历"),
]
for i,(accent,title,dir_,schools,exam,tip) in enumerate(routes):
    y = 1.5 + i*1.7
    _round_box(sl, 1.0, y, 11.3, 1.5, WHITE)
    _box(sl, 1.0, y, 0.1, 1.5, accent)
    _tb(sl, 1.35, y+0.12, 2.5, 0.4, title, fs=18, bold=True, color=accent)
    _tb(sl, 1.35, y+0.55, 5, 0.3, f"方向：{dir_}", fs=13, bold=True, color=DARK)
    _tb(sl, 1.35, y+0.9, 4, 0.25, f"目标院校：{schools}", fs=11, color=BLUE)
    _tb(sl, 6.5, y+0.25, 3, 0.3, f"考试科目：{exam}", fs=13, bold=True, color=RED)
    _tb(sl, 6.5, y+0.7, 5.5, 0.4, f"💡 {tip}", fs=11, color=GRAY)

exam_tl = ["大二暑假\n确定方向", "大三上\n数学/英语", "大三寒假\n真题集训", "大三下\n专业课", "大四上\n冲刺模考"]
for i,item in enumerate(exam_tl):
    x = 1.0 + i*2.4
    _circle(sl, x+0.7, 6.55, 0.22, RED)
    if i<4: _box(sl, x+0.95, 6.64, 2.2, 0.04, RED)
    _tb(sl, x, 6.85, 2.2, 0.3, item, fs=9, bold=True, color=DARK, align=PP_ALIGN.CENTER)

_footer(sl, 16)


# ============================================================
# SLIDE 17 — PATH 3: GOV
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
_page_head(sl, "路径三：公务员与事业单位 —— 公共服务赛道")

gov_jobs = [
    ("国考", "住建部 / 发改委", "高度对口", "1:100+"),
    ("省考", "省住建厅 / 自然资源厅", "高度对口", "中等偏高"),
    ("市考", "市住建局 / 区建设局", "对口", "中等"),
    ("事业单位", "工程质量监督站 / 公共资源交易中心", "对口", "中等"),
    ("国企编制", "城投公司 / 保障房建设公司", "高度对口", "中等偏低"),
]
for i,(level,dept,fit,comp) in enumerate(gov_jobs):
    y = 1.5 + i*0.75
    _round_box(sl, 1.0, y, 11.3, 0.6, WHITE if i%2==0 else LIGHT)
    _tb(sl, 1.2, y+0.1, 2, 0.35, level, fs=15, bold=True, color=RED)
    _tb(sl, 3.3, y+0.1, 4.5, 0.35, dept, fs=13, color=DARK)
    _tb(sl, 7.8, y+0.1, 2, 0.35, fit, fs=13, bold=True, color=GREEN)
    _tb(sl, 10.5, y+0.1, 2, 0.35, comp, fs=12, color=GRAY)

# advantage box
_round_box(sl, 1.0, 5.4, 11.3, 1.3, GREEN)
_tb(sl, 1.3, 5.5, 4, 0.4, "西建大考公核心优势", fs=17, bold=True, color=WHITE)
adv = [
    "住建系统对西建大认可度较高，\"建筑老八校\"品牌在体制内有辨识度",
    "建筑类岗位限制专业，竞争比远低于\"三不限\"通用岗位",
    "校友网络广泛分布于西北地区住建系统，可提前了解招录动态",
]
tf = _rbox(sl, 1.3, 5.95, 10.5, 0.7)
for a in adv:
    _p(tf, f"✦ {a}", fs=12, color=WHITE, sa=Pt(4))

_tb(sl, 1.0, 6.85, 11, 0.3,
    "👤 适合人群：追求稳定 · 对公共政策感兴趣 · 重视工作生活平衡 · 愿意长期在体制内发展",
    fs=13, bold=True, color=DARK)
_footer(sl, 17)


# ============================================================
# SLIDE 18 — GROUP & OPINIONS
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
_page_head(sl, "小组分工与成员观点")

# roles - just simple grid
_box(sl, 1.0, 1.5, 11.3, 2.6, LIGHT)
_tb(sl, 1.3, 1.6, 3, 0.35, "📋 小组分工", fs=15, bold=True, color=DARK)

roles = [
    ("成员A（组长）", "整体框架 & 第一二部分", "专业现状与趋势分析"),
    ("成员B", "第三部分：学习规划", "四年时间轴 & 技能提升路径"),
    ("成员C", "第四部分：职业方向", "三路径解析 & 数据搜集"),
    ("成员D", "设计美化 & 数据可视化", "PPT 排版 & 图表制作"),
    ("成员E", "资料搜集 & 演讲备注", "课程数据 & 就业数据整理"),
    ("成员F", "审校 & 模拟演练", "内容审核 & 演讲演练"),
]
for i,(name,module,contribution) in enumerate(roles):
    col = i % 3
    row = i // 3
    x = 1.3 + col*3.7
    y = 2.1 + row*0.95
    _round_box(sl, x, y, 3.4, 0.78, WHITE)
    _tb(sl, x+0.15, y+0.05, 3.1, 0.28, name, fs=13, bold=True, color=RED)
    _tb(sl, x+0.15, y+0.32, 3.1, 0.2, module, fs=10, color=DARK)
    _tb(sl, x+0.15, y+0.52, 3.1, 0.2, contribution, fs=9, color=GRAY)

# opinions
_tb(sl, 1.0, 4.3, 4, 0.35, "💬 成员个人观点摘录", fs=15, bold=True, color=DARK)
opinions = [
    ("🟡 务实就业派", "想进建筑行业做管理，大三去陕建项目实习，从管培生做起，目标5-8年做到项目经理。"),
    ("🟢 考研深造派", "大三备考同济大学管科，读研专攻工程项目管理，毕业选择面会宽很多。"),
    ("🟠 考公稳定派", "参加国考或省考，目标住建局。\"建筑老八校\"背景在专业限制岗位上有优势。"),
    ("🔵 跨界转型派", "想尝试建筑科技方向。自学Python，目标广联达做商务数据分析。"),
    ("🟣 国际工程派", "英语底子好，想去中建海外或咨询外企，走国际商务方向。"),
    ("🔴 探索观望派", "还没完全确定，大二多参加商赛，大三实习后再定，先保证不挂科+四六级。"),
]
for i,(label,opinion) in enumerate(opinions):
    col = i % 3
    row = i // 3
    x = 1.0 + col*4.0
    y = 4.75 + row*1.15
    _round_box(sl, x, y, 3.7, 1.0, WHITE)
    _box(sl, x, y, 3.7, 0.06, RED)
    _tb(sl, x+0.15, y+0.12, 3.4, 0.3, label, fs=12, bold=True, color=DARK)
    _tb(sl, x+0.15, y+0.42, 3.4, 0.5, f"\"{opinion}\"", fs=10, color=DARK)
_footer(sl, 18)


# ============================================================
# SLIDE 19 — END
# ============================================================
sl = prs.slides.add_slide(prs.slide_layouts[6])
_box(sl, 0, 0, 13.333, 7.5, DARK)
_box(sl, 1.5, 2.0, 0.12, 4.0, RED)

_tb(sl, 2.2, 2.0, 9, 1.0, "规划未来 · 行则将至", fs=44, bold=True, color=WHITE)

lines = [
    ("今天的规划，是明天的底气。", 18, WHITE),
    ("", 4, WHITE),
    ("工商管理是一个\"上限很高\"的专业 —— 它不会限定你的职业边界，", 14, GRAY),
    ("但也不会替你铺好每一步路。", 14, GRAY),
    ("", 4, WHITE),
    ("在西安建筑科技大学管理学院，我们拥有：", 14, WHITE),
    ("   ▶  60 年建筑经济管理的历史积淀", 14, WHITE),
    ("   ▶  国家级一流专业的平台支撑", 14, WHITE),
    ("   ▶  \"建筑 + 管理 + 数智\" 的独特定位", 14, WHITE),
    ("", 4, WHITE),
    ("无论选择就业、升学还是考公，提前 2-3 年规划，", 16, WHITE),
    ("一步一个脚印执行，就是我们最大的竞争力。", 16, WHITE),
]
tf = _rbox(sl, 2.2, 3.3, 9.5, 3.0)
for text,fs,color in lines:
    _p(tf, text, fs=fs, color=color, sa=Pt(4))

_box(sl, 2.2, 6.3, 4, 0.05, RED)
_tb(sl, 2.2, 6.5, 9, 0.55, "感谢聆听！", fs=24, bold=True, color=WHITE)
_tb(sl, 2.2, 6.9, 9, 0.3, "西安建筑科技大学管理学院 · 工商管理专业 · 宿舍小组 · 2026年5月", fs=11, color=GRAY)


# ============================================================
# SAVE
# ============================================================
os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print(f"PPT saved: {OUT}")
print(f"Slides: {len(prs.slides)}")
