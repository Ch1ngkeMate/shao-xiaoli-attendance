"""
生成「如何使用 Claude Code — 从入门到精通」PPT
设计风格：深蓝色现代风格，14页
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn
import os

# ============================================================
# 颜色常量
# ============================================================
BG_DARK      = RGBColor(0x1A, 0x1A, 0x2E)
BG_CARD      = RGBColor(0x16, 0x21, 0x3E)
ACCENT_RED   = RGBColor(0xE9, 0x45, 0x60)
ACCENT_GOLD  = RGBColor(0xF0, 0xA5, 0x00)
ACCENT_BLUE  = RGBColor(0x53, 0x7F, 0xE7)
ACCENT_GREEN = RGBColor(0x2E, 0xCC, 0x71)
ACCENT_PURPLE= RGBColor(0xA5, 0x5E, 0xE0)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xB0, 0xB0, 0xB0)
DARK_GRAY    = RGBColor(0x8A, 0x8A, 0x8A)
MID_BLUE     = RGBColor(0x0F, 0x34, 0x60)
BORDER_GRAY  = RGBColor(0x33, 0x33, 0x55)

# 16:9 幻灯片
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT_TITLE = 'Microsoft YaHei'
FONT_BODY  = 'Microsoft YaHei'
FONT_EN    = 'Calibri'

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

# ============================================================
# 辅助函数
# ============================================================

def add_blank_slide():
    """添加空白幻灯片"""
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)

def set_bg(slide, color=BG_DARK):
    """设置幻灯片背景色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None, border_color=None, border_width=None, corner_radius=None):
    """添加矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if corner_radius else MSO_SHAPE.RECTANGLE,
                                   left, top, width, height)
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.fill.solid()
        if border_width:
            shape.line.width = border_width
    return shape

def add_textbox(slide, left, top, width, height, text="", font_size=Pt(14),
                font_color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                font_name=FONT_BODY, line_spacing=1.2, anchor=MSO_ANCHOR.TOP):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    # 设置垂直锚点
    txBox.text_frame.auto_size = None
    try:
        tf.paragraphs[0].alignment = alignment
    except:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.line_spacing = line_spacing
    # 设置东亚字体
    for run in p.runs:
        rPr = run._r.get_or_add_rPr()
        rPr.set(qn('a:eaTypeface'), font_name if font_name == FONT_TITLE else FONT_BODY)
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines, default_size=Pt(13),
                          default_color=WHITE, default_bold=False, alignment=PP_ALIGN.LEFT,
                          font_name=FONT_BODY, line_spacing=1.3):
    """添加多行文本框，lines = [(text, size, color, bold), ...]"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if isinstance(line_data, str):
            text, size, color, bold = line_data, default_size, default_color, default_bold
        else:
            text = line_data[0]
            size = line_data[1] if len(line_data) > 1 else default_size
            color = line_data[2] if len(line_data) > 2 else default_color
            bold = line_data[3] if len(line_data) > 3 else default_bold
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = size
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.line_spacing = line_spacing
        for run in p.runs:
            rPr = run._r.get_or_add_rPr()
            rPr.set(qn('a:eaTypeface'), font_name)
    return txBox

def add_decorative_line(slide, left, top, width, color=ACCENT_RED, height=Pt(2)):
    """添加装饰线"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_page_number(slide, num):
    """右下角页码"""
    add_textbox(slide, Inches(12.3), Inches(7.05), Inches(0.8), Inches(0.35),
                str(num), Pt(10), DARK_GRAY, alignment=PP_ALIGN.RIGHT)

def add_section_indicator(slide, color):
    """左侧章节色块指示器"""
    add_rect(slide, Inches(0), Inches(0), Inches(0.08), Inches(7.5), fill_color=color)

def add_top_accent(slide, color=ACCENT_RED):
    """顶部装饰线"""
    add_rect(slide, Inches(0.6), Inches(0.25), Inches(12.1), Pt(2), fill_color=color)

def add_slide_title(slide, title, subtitle=None, title_color=WHITE):
    """统一标题区域"""
    add_top_accent(slide)
    add_textbox(slide, Inches(0.8), Inches(0.45), Inches(11.5), Inches(0.55),
                title, Pt(28), title_color, bold=True, font_name=FONT_TITLE)
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(0.95), Inches(11.5), Inches(0.35),
                    subtitle, Pt(13), LIGHT_GRAY, font_name=FONT_BODY)

def add_card(slide, left, top, width, height, title, desc_lines, icon_text=None,
             accent_color=ACCENT_RED, title_size=Pt(15), desc_size=Pt(11)):
    """添加卡片"""
    add_rect(slide, left, top, width, height, fill_color=BG_CARD, corner_radius=True)
    # 顶部色条
    add_rect(slide, left + Inches(0.01), top, width - Inches(0.02), Pt(3), fill_color=accent_color)
    # 图标(如果有)
    y_offset = top + Inches(0.2)
    if icon_text:
        add_textbox(slide, left + Inches(0.2), y_offset, width - Inches(0.4), Inches(0.35),
                    icon_text, Pt(20), accent_color, bold=True, font_name=FONT_EN)
        y_offset += Inches(0.35)
    # 标题
    add_textbox(slide, left + Inches(0.2), y_offset, width - Inches(0.4), Inches(0.3),
                title, title_size, WHITE, bold=True, font_name=FONT_TITLE)
    y_offset += Inches(0.35)
    # 描述
    for line in desc_lines:
        add_textbox(slide, left + Inches(0.2), y_offset, width - Inches(0.4), Inches(0.25),
                    line, desc_size, LIGHT_GRAY, font_name=FONT_BODY)
        y_offset += Inches(0.22)

# ============================================================
# 第1页：封面
# ============================================================
slide1 = add_blank_slide()
set_bg(slide1, BG_DARK)
# 装饰几何图形 - 大圆
circle = slide1.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(-1.5), Inches(5.5), Inches(5.5))
circle.fill.background()
circle.line.color.rgb = ACCENT_RED
circle.line.width = Pt(1)
circle.line.dash_style = 2  # dash
# 小方块装饰
add_rect(slide1, Inches(1.0), Inches(5.8), Inches(0.5), Inches(0.5), fill_color=ACCENT_RED, corner_radius=False)
add_rect(slide1, Inches(1.7), Inches(6.0), Inches(0.4), Inches(0.4), fill_color=None, border_color=ACCENT_GOLD, border_width=Pt(1.5), corner_radius=False)
add_rect(slide1, Inches(11.0), Inches(6.2), Inches(0.6), Inches(0.6), fill_color=ACCENT_BLUE, corner_radius=False)
# 主标题
add_textbox(slide1, Inches(1.0), Inches(1.8), Inches(10.0), Inches(1.2),
            "如何使用 Claude Code", Pt(52), WHITE, bold=True, font_name=FONT_TITLE)
# 红线装饰
add_rect(slide1, Inches(1.0), Inches(3.0), Inches(2.5), Pt(4), fill_color=ACCENT_RED)
# 副标题
add_textbox(slide1, Inches(1.0), Inches(3.3), Inches(8.0), Inches(0.6),
            "AI 编程助手从入门到精通", Pt(22), LIGHT_GRAY, font_name=FONT_TITLE)
# 描述
add_textbox(slide1, Inches(1.0), Inches(4.1), Inches(8.0), Inches(0.8),
            "全面了解 Claude Code 的功能、技巧与最佳实践，\n让 AI 编程助手成为你的高效开发伙伴。",
            Pt(14), DARK_GRAY, font_name=FONT_BODY)
# 日期
add_textbox(slide1, Inches(1.0), Inches(5.2), Inches(4.0), Inches(0.3),
            "2026 年 5 月", Pt(13), DARK_GRAY, font_name=FONT_EN)
add_page_number(slide1, 1)

# ============================================================
# 第2页：目录
# ============================================================
slide2 = add_blank_slide()
set_bg(slide2, BG_DARK)
add_slide_title(slide2, "目录", "CONTENTS")

chapters = [
    ("01", "认识 Claude", "版本介绍 · 核心能力 · 适用场景", ACCENT_RED),
    ("02", "Claude Code CLI", "安装配置 · 环境搭建 · 基本使用", ACCENT_BLUE),
    ("03", "核心功能详解", "文件操作 · 代码生成 · 审查重构", ACCENT_GOLD),
    ("04", "Git 工作流", "版本控制 · 分支管理 · Pull Request", ACCENT_GREEN),
    ("05", "高级技巧", "Agent 并行 · Memory 系统 · 自定义配置", ACCENT_PURPLE),
    ("06", "最佳实践", "高效沟通 · 场景案例 · 资源汇总", ACCENT_RED),
]

card_w = Inches(3.7)
card_h = Inches(2.2)
gap_x = Inches(0.35)
gap_y = Inches(0.3)
start_x = Inches(0.8)
start_y = Inches(1.7)

for i, (num, title, desc, color) in enumerate(chapters):
    col = i % 3
    row = i // 3
    x = start_x + col * (card_w + gap_x)
    y = start_y + row * (card_h + gap_y)
    add_rect(slide2, x, y, card_w, card_h, fill_color=BG_CARD, corner_radius=True)
    add_rect(slide2, x + Inches(0.01), y, card_w - Inches(0.02), Pt(3), fill_color=color)
    add_textbox(slide2, x + Inches(0.25), y + Inches(0.2), Inches(0.6), Inches(0.35),
                num, Pt(26), color, bold=True, font_name=FONT_EN)
    add_textbox(slide2, x + Inches(0.25), y + Inches(0.7), card_w - Inches(0.5), Inches(0.35),
                title, Pt(17), WHITE, bold=True, font_name=FONT_TITLE)
    add_textbox(slide2, x + Inches(0.25), y + Inches(1.3), card_w - Inches(0.5), Inches(0.6),
                desc, Pt(11), LIGHT_GRAY, font_name=FONT_BODY)
add_page_number(slide2, 2)

# ============================================================
# 第3页：认识 Claude
# ============================================================
slide3 = add_blank_slide()
set_bg(slide3, BG_DARK)
add_slide_title(slide3, "认识 Claude", "什么是 Claude？它为什么与众不同？")
add_section_indicator(slide3, ACCENT_RED)

# 左侧介绍
add_textbox(slide3, Inches(0.8), Inches(1.6), Inches(5.8), Inches(0.3),
            "Claude 是什么？", Pt(20), ACCENT_RED, bold=True, font_name=FONT_TITLE)
add_multiline_textbox(slide3, Inches(0.8), Inches(2.1), Inches(5.8), Inches(3.5), [
    ("Claude 是由 Anthropic 开发的先进大语言模型（LLM），", Pt(14), WHITE, False),
    ("专注于安全、可靠、有帮助的 AI 交互体验。", Pt(14), WHITE, False),
    ("", Pt(8), WHITE, False),
    ("核心特点：", Pt(16), ACCENT_GOLD, True),
    ("", Pt(6), WHITE, False),
    ("◆ 强大的代码理解与生成能力", Pt(13), WHITE, False),
    ("◆ 支持多文件编辑、重构与调试", Pt(13), WHITE, False),
    ("◆ 内置安全机制，减少有害输出", Pt(13), WHITE, False),
    ("◆ 可通过 Claude Code CLI 在终端中直接使用", Pt(13), WHITE, False),
    ("◆ 支持 Agent 模式，自动完成复杂任务", Pt(13), WHITE, False),
    ("◆ Memory 系统实现跨会话持久化记忆", Pt(13), WHITE, False),
])

# 右侧关键数据卡片
data_cards = [
    ("模型版本", "3 种", "Opus · Sonnet · Haiku", ACCENT_RED),
    ("上下文窗口", "200K tokens", "可处理大型代码库", ACCENT_BLUE),
    ("支持语言", "100+", "主流编程语言全覆盖", ACCENT_GREEN),
]
for i, (label, value, desc, color) in enumerate(data_cards):
    y = Inches(1.6) + i * Inches(1.7)
    add_rect(slide3, Inches(7.5), y, Inches(5.0), Inches(1.4), fill_color=BG_CARD, corner_radius=True)
    add_rect(slide3, Inches(7.52), y, Inches(0.06), Inches(1.4), fill_color=color)
    add_textbox(slide3, Inches(8.0), y + Inches(0.15), Inches(4.2), Inches(0.25),
                label, Pt(11), LIGHT_GRAY, font_name=FONT_BODY)
    add_textbox(slide3, Inches(8.0), y + Inches(0.4), Inches(4.2), Inches(0.45),
                value, Pt(28), color, bold=True, font_name=FONT_EN)
    add_textbox(slide3, Inches(8.0), y + Inches(0.9), Inches(4.2), Inches(0.25),
                desc, Pt(12), LIGHT_GRAY, font_name=FONT_BODY)
add_page_number(slide3, 3)

# ============================================================
# 第4页：模型家族
# ============================================================
slide4 = add_blank_slide()
set_bg(slide4, BG_DARK)
add_slide_title(slide4, "模型家族", "选择适合你任务的 Claude 模型")
add_section_indicator(slide4, ACCENT_BLUE)

models = [
    ("Claude Opus", "最强性能", ACCENT_RED,
     ["最强大的推理与创作能力", "适合复杂架构设计、代码审查", "深度分析与决策支持", "处理高难度技术问题"]),
    ("Claude Sonnet", "最佳平衡", ACCENT_BLUE,
     ["性能与速度的理想平衡", "日常编程任务的优选", "快速代码生成与调试", "PR 审查与文档编写"]),
    ("Claude Haiku", "极速响应", ACCENT_GREEN,
     ["毫秒级响应速度", "适合简单查询与交互", "轻量级代码补全", "实时对话场景"]),
]

card_w4 = Inches(3.7)
gap4 = Inches(0.35)
start_x4 = Inches(0.8)

for i, (name, tagline, color, features) in enumerate(models):
    x = start_x4 + i * (card_w4 + gap4)
    y = Inches(1.6)
    # 卡片背景
    add_rect(slide4, x, y, card_w4, Inches(5.0), fill_color=BG_CARD, corner_radius=True)
    add_rect(slide4, x + Inches(0.01), y, card_w4 - Inches(0.02), Pt(4), fill_color=color)
    # 模型名
    add_textbox(slide4, x + Inches(0.3), y + Inches(0.3), card_w4 - Inches(0.6), Inches(0.4),
                name, Pt(24), WHITE, bold=True, font_name=FONT_EN)
    # 标签
    tag = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x + Inches(0.3), y + Inches(0.8), Inches(1.5), Inches(0.3))
    tag.fill.solid()
    tag.fill.fore_color.rgb = color
    tag.line.fill.background()
    tf = tag.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = tagline
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_BODY
    p.alignment = PP_ALIGN.CENTER
    # 特性
    fy = y + Inches(1.4)
    for feat in features:
        add_textbox(slide4, x + Inches(0.3), fy, card_w4 - Inches(0.6), Inches(0.3),
                    "✓  " + feat, Pt(12), LIGHT_GRAY, font_name=FONT_BODY)
        fy += Inches(0.38)

add_page_number(slide4, 4)

# ============================================================
# 第5页：Claude Code CLI
# ============================================================
slide5 = add_blank_slide()
set_bg(slide5, BG_DARK)
add_slide_title(slide5, "Claude Code CLI", "安装、配置与环境搭建")
add_section_indicator(slide5, ACCENT_BLUE)

# 左侧安装步骤
add_textbox(slide5, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.3),
            "安装步骤", Pt(20), ACCENT_RED, bold=True, font_name=FONT_TITLE)
steps = [
    ("1", "安装 Node.js 18+", "从 nodejs.org 下载并安装"),
    ("2", "全局安装 CLI", "npm install -g @anthropic-ai/claude-code"),
    ("3", "获取 API Key", "在 console.anthropic.com 创建密钥"),
    ("4", "设置环境变量", "set ANTHROPIC_API_KEY=your-key"),
    ("5", "启动 Claude Code", "在终端输入 claude 开始使用"),
]
for i, (num, title, desc) in enumerate(steps):
    y = Inches(2.15) + i * Inches(0.9)
    add_rect(slide5, Inches(0.8), y, Inches(5.5), Inches(0.75), fill_color=BG_CARD, corner_radius=True)
    # 编号圆圈
    circle = slide5.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y + Inches(0.15), Inches(0.45), Inches(0.45))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT_RED
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_EN
    add_textbox(slide5, Inches(1.7), y + Inches(0.08), Inches(4.3), Inches(0.25),
                title, Pt(13), WHITE, bold=True, font_name=FONT_TITLE)
    add_textbox(slide5, Inches(1.7), y + Inches(0.38), Inches(4.3), Inches(0.22),
                desc, Pt(11), LIGHT_GRAY, font_name=FONT_EN)

# 右侧提示
add_textbox(slide5, Inches(7.2), Inches(1.6), Inches(5.3), Inches(0.3),
            "配置建议", Pt(20), ACCENT_GOLD, bold=True, font_name=FONT_TITLE)
tips_text = [
    "▸ 使用 CLAUDE.md 文件定义项目级指令",
    "▸ 配置 settings.json 管理权限与行为",
    "▸ 利用 Memory 系统保存偏好与经验",
    "▸ 设置 /config 调整主题、模型等参数",
    "▸ 通过 hooks 自动化重复性操作",
    "▸ 多项目切换时注意工作目录上下文",
]
for i, tip in enumerate(tips_text):
    y = Inches(2.15) + i * Inches(0.55)
    add_textbox(slide5, Inches(7.2), y, Inches(5.3), Inches(0.45),
                tip, Pt(13), WHITE, font_name=FONT_BODY)

add_page_number(slide5, 5)

# ============================================================
# 第6页：核心功能详解
# ============================================================
slide6 = add_blank_slide()
set_bg(slide6, BG_DARK)
add_slide_title(slide6, "核心功能详解", "Claude Code 能为你做什么？")
add_section_indicator(slide6, ACCENT_GOLD)

features = [
    ("📄", "文件操作", "读取、搜索、编辑、\n创建文件，支持批量处理", ACCENT_RED),
    ("⚡", "代码生成", "从零创建项目模块，\n理解需求自动编码", ACCENT_BLUE),
    ("🔍", "代码审查", "Review 代码变更，\n发现潜在问题与漏洞", ACCENT_GOLD),
    ("🔧", "Bug 修复", "定位并修复 Bug，\n提供详细的修复方案", ACCENT_GREEN),
    ("♻️", "代码重构", "优化代码结构，\n提升可读性与性能", ACCENT_PURPLE),
    ("🔗", "Git 集成", "管理分支、提交、\n创建 Pull Request", ACCENT_RED),
]

fw = Inches(3.7)
fh = Inches(2.2)
fgap_x = Inches(0.35)
fgap_y = Inches(0.3)
fx_start = Inches(0.8)
fy_start = Inches(1.6)

for i, (icon, title, desc, color) in enumerate(features):
    col = i % 3
    row = i // 3
    x = fx_start + col * (fw + fgap_x)
    y = fy_start + row * (fh + fgap_y)
    add_rect(slide6, x, y, fw, fh, fill_color=BG_CARD, corner_radius=True)
    add_rect(slide6, x, y, fw, Pt(3), fill_color=color)
    add_textbox(slide6, x + Inches(0.25), y + Inches(0.25), Inches(0.5), Inches(0.4),
                icon, Pt(24), font_name=FONT_EN)
    add_textbox(slide6, x + Inches(0.25), y + Inches(0.75), fw - Inches(0.5), Inches(0.3),
                title, Pt(17), WHITE, bold=True, font_name=FONT_TITLE)
    add_textbox(slide6, x + Inches(0.25), y + Inches(1.2), fw - Inches(0.5), Inches(0.8),
                desc, Pt(12), LIGHT_GRAY, font_name=FONT_BODY)
add_page_number(slide6, 6)

# ============================================================
# 第7页：文件操作详解
# ============================================================
slide7 = add_blank_slide()
set_bg(slide7, BG_DARK)
add_slide_title(slide7, "文件操作详解", "高效管理代码文件的核心能力")
add_section_indicator(slide7, ACCENT_RED)

ops = [
    ("Read", "读取文件", "查看任意文件内容，支持分页、\n行号显示，图片与 PDF 也能解析。", ACCENT_BLUE,
     ["read_file 或 Read 工具", "自动显示行号", "支持图片/PDF 解析", "分页读取大文件"]),
    ("Edit", "精确编辑", "通过字符串匹配进行精确替换，\n避免重写整个文件。", ACCENT_GREEN,
     ["精确字符串匹配替换", "支持 replace_all 批量替换", "保持原有缩进格式", "编辑失败自动回滚"]),
    ("Write", "写入文件", "创建新文件或完整覆写，\n适用于从零生成代码。", ACCENT_GOLD,
     ["创建新文件或覆写", "自动处理路径", "支持大文件生成", "配合 Glob/Grep 使用"]),
]

for i, (label, title, desc, color, details) in enumerate(ops):
    x = Inches(0.8) + i * Inches(4.1)
    y = Inches(1.6)
    add_rect(slide7, x, y, Inches(3.8), Inches(5.2), fill_color=BG_CARD, corner_radius=True)
    add_rect(slide7, x, y, Inches(3.8), Pt(3), fill_color=color)
    # 标签
    tag = slide7.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   x + Inches(0.25), y + Inches(0.25), Inches(1.0), Inches(0.3))
    tag.fill.solid()
    tag.fill.fore_color.rgb = color
    tag.line.fill.background()
    tf = tag.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = FONT_EN
    p.alignment = PP_ALIGN.CENTER
    # 标题
    add_textbox(slide7, x + Inches(0.25), y + Inches(0.75), Inches(3.3), Inches(0.3),
                title, Pt(18), WHITE, bold=True, font_name=FONT_TITLE)
    # 描述
    add_textbox(slide7, x + Inches(0.25), y + Inches(1.2), Inches(3.3), Inches(0.7),
                desc, Pt(12), LIGHT_GRAY, font_name=FONT_BODY)
    # 详情
    dy = y + Inches(2.2)
    for d in details:
        add_textbox(slide7, x + Inches(0.25), dy, Inches(3.3), Inches(0.28),
                    "• " + d, Pt(11), LIGHT_GRAY, font_name=FONT_BODY)
        dy += Inches(0.32)
add_page_number(slide7, 7)

# ============================================================
# 第8页：代码生成
# ============================================================
slide8 = add_blank_slide()
set_bg(slide8, BG_DARK)
add_slide_title(slide8, "代码生成", "从零到一，快速构建项目")
add_section_indicator(slide8, ACCENT_BLUE)

# 左侧流程
add_textbox(slide8, Inches(0.8), Inches(1.6), Inches(5.5), Inches(0.3),
            "工作流程", Pt(20), ACCENT_RED, bold=True, font_name=FONT_TITLE)
flow_steps = [
    ("需求描述", "清晰描述你想要的功能或项目结构"),
    ("Plan Mode", "Claude 先制定计划，确认后再编码"),
    ("逐文件生成", "按计划逐步生成代码文件"),
    ("测试验证", "自动运行测试，确保功能正确"),
    ("迭代优化", "根据反馈进行调整和改进"),
]
for i, (step, desc) in enumerate(flow_steps):
    y = Inches(2.2) + i * Inches(0.85)
    # 连接线
    if i < len(flow_steps) - 1:
        add_rect(slide8, Inches(1.25), y + Inches(0.55), Pt(2), Inches(0.3), fill_color=ACCENT_RED)
    # 圆点
    dot = slide8.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.1), y + Inches(0.1), Inches(0.3), Inches(0.3))
    dot.fill.solid()
    dot.fill.fore_color.rgb = ACCENT_RED
    dot.line.fill.background()
    add_textbox(slide8, Inches(1.7), y, Inches(4.3), Inches(0.28),
                step, Pt(14), WHITE, bold=True, font_name=FONT_TITLE)
    add_textbox(slide8, Inches(1.7), y + Inches(0.3), Inches(4.3), Inches(0.25),
                desc, Pt(11), LIGHT_GRAY, font_name=FONT_BODY)

# 右侧示例
add_textbox(slide8, Inches(7.2), Inches(1.6), Inches(5.3), Inches(0.3),
            "实用技巧", Pt(20), ACCENT_GOLD, bold=True, font_name=FONT_TITLE)
tips = [
    "生成前的关键准备：",
    "",
    "① 提供清晰的技术栈说明",
    "   （如 Next.js + TypeScript + Prisma）",
    "② 描述期望的文件结构",
    "③ 指定代码风格与命名规范",
    "④ 说明关键业务逻辑与边界条件",
    "⑤ 使用 Plan Mode 降低返工成本",
    "⑥ 对生成代码执行 Type Check",
]
for i, tip in enumerate(tips):
    y = Inches(2.2) + i * Inches(0.42)
    color = ACCENT_GOLD if tip.startswith("生成") else (WHITE if tip.startswith("①") or tip.startswith("②") or tip.startswith("③") or tip.startswith("④") or tip.startswith("⑤") or tip.startswith("⑥") else LIGHT_GRAY)
    bold = tip.startswith("生成") or tip.startswith("①") or tip.startswith("②") or tip.startswith("③") or tip.startswith("④") or tip.startswith("⑤") or tip.startswith("⑥")
    add_textbox(slide8, Inches(7.2), y, Inches(5.3), Inches(0.35),
                tip, Pt(12), color, bold=bold, font_name=FONT_BODY if not bold else FONT_TITLE)
add_page_number(slide8, 8)

# ============================================================
# 第9页：审查与重构
# ============================================================
slide9 = add_blank_slide()
set_bg(slide9, BG_DARK)
add_slide_title(slide9, "代码审查与重构", "提升代码质量的利器")
add_section_indicator(slide9, ACCENT_GREEN)

# 三个并排卡片
review_items = [
    ("PR Review", ACCENT_BLUE,
     ["自动检查代码变更", "发现潜在 Bug 与逻辑错误", "评估安全风险", "提供改进建议", "检查代码风格一致性"]),
    ("代码重构", ACCENT_GREEN,
     ["优化函数与类结构", "消除重复代码", "提升可读性", "改进命名规范", "应用设计模式"]),
    ("安全审查", ACCENT_RED,
     ["检测 SQL 注入风险", "识别 XSS 漏洞", "审查认证逻辑", "检查敏感数据暴露", "验证权限控制"]),
]

for i, (title, color, items) in enumerate(review_items):
    x = Inches(0.8) + i * Inches(4.1)
    y = Inches(1.6)
    add_rect(slide9, x, y, Inches(3.8), Inches(5.0), fill_color=BG_CARD, corner_radius=True)
    add_rect(slide9, x, y, Inches(3.8), Pt(3), fill_color=color)
    add_textbox(slide9, x + Inches(0.25), y + Inches(0.25), Inches(3.3), Inches(0.35),
                title, Pt(20), color, bold=True, font_name=FONT_EN)
    add_rect(slide9, x + Inches(0.25), y + Inches(0.7), Inches(1.5), Pt(2), fill_color=color)
    dy = y + Inches(1.0)
    for item in items:
        add_textbox(slide9, x + Inches(0.25), dy, Inches(3.3), Inches(0.28),
                    "▹ " + item, Pt(12), LIGHT_GRAY, font_name=FONT_BODY)
        dy += Inches(0.38)
add_page_number(slide9, 9)

# ============================================================
# 第10页：Git 工作流
# ============================================================
slide10 = add_blank_slide()
set_bg(slide10, BG_DARK)
add_slide_title(slide10, "Git 工作流集成", "Claude Code 如何助力版本控制")
add_section_indicator(slide10, ACCENT_GREEN)

# 左侧工作流图
add_textbox(slide10, Inches(0.8), Inches(1.6), Inches(6.0), Inches(0.3),
            "完整的 Git 操作流程", Pt(20), ACCENT_RED, bold=True, font_name=FONT_TITLE)

git_steps = [
    ("查看状态", "git status — 了解当前变更"),
    ("暂存文件", "git add — 精确选择要提交的文件"),
    ("生成提交信息", "Claude 自动分析变更并生成规范的 commit message"),
    ("创建提交", "git commit — 安全提交，自动处理 hook"),
    ("管理分支", "git branch / checkout — 功能分支管理"),
    ("创建 PR", "gh pr create — Claude 自动生成标题与描述"),
]

for i, (step, desc) in enumerate(git_steps):
    y = Inches(2.2) + i * Inches(0.75)
    add_rect(slide10, Inches(0.8), y, Inches(6.0), Inches(0.6), fill_color=BG_CARD, corner_radius=True)
    add_textbox(slide10, Inches(1.0), y + Inches(0.08), Inches(1.8), Inches(0.22),
                step, Pt(12), ACCENT_GOLD, bold=True, font_name=FONT_TITLE)
    add_textbox(slide10, Inches(1.0), y + Inches(0.32), Inches(5.5), Inches(0.22),
                desc, Pt(11), LIGHT_GRAY, font_name=FONT_EN)

# 右侧安全提示
add_textbox(slide10, Inches(7.5), Inches(1.6), Inches(5.0), Inches(0.3),
            "Git 安全协议", Pt(20), ACCENT_RED, bold=True, font_name=FONT_TITLE)

safety_items = [
    "⚠ 永远不要跳过 Git Hooks",
    "⚠ 不要修改 Git 配置",
    "⚠ 推送前确认分支正确",
    "⚠ Force push 需用户明确授权",
    "⚠ 不做 git reset --hard 等破坏性操作",
    "⚠ 敏感文件（.env）自动排除",
    "✓ 始终创建新提交而非 amend",
    "✓ 使用规范的 commit message",
    "✓ PR 自动生成测试计划",
]
for i, item in enumerate(safety_items):
    y = Inches(2.2) + i * Inches(0.45)
    color = ACCENT_RED if item.startswith("⚠") else ACCENT_GREEN
    add_textbox(slide10, Inches(7.5), y, Inches(5.0), Inches(0.35),
                item, Pt(12), color, bold=False, font_name=FONT_BODY)
add_page_number(slide10, 10)

# ============================================================
# 第11页：高级技巧
# ============================================================
slide11 = add_blank_slide()
set_bg(slide11, BG_DARK)
add_slide_title(slide11, "高级技巧", "解锁 Claude Code 的全部潜力")
add_section_indicator(slide11, ACCENT_PURPLE)

advanced = [
    ("Agent 并行", ACCENT_RED, Inches(0.8), Inches(1.6),
     ["将独立任务分发到多个子 Agent", "并行执行提高效率", "支持 Explore / Plan / Review 等类型",
      "结果自动汇总到主会话", "适合大型代码库探索与审查"]),
    ("Memory 系统", ACCENT_BLUE, Inches(4.95), Inches(1.6),
     ["跨会话持久化记忆", "User / Project / Feedback 等类型",
      "自动保存用户偏好与经验", "下次对话自动加载相关记忆",
      "可随时更新或删除记忆"]),
    ("自定义配置", ACCENT_GREEN, Inches(9.1), Inches(1.6),
     [".claude/settings.json 管理权限", "CLAUDE.md 定义项目级指令",
      "Hooks 自动触发操作", "Worktree 隔离工作环境",
      "Cron 定时任务调度"]),
]

for title, color, x, y, items in advanced:
    add_rect(slide11, x, y, Inches(3.85), Inches(3.8), fill_color=BG_CARD, corner_radius=True)
    add_rect(slide11, x, y, Inches(3.85), Pt(3), fill_color=color)
    add_textbox(slide11, x + Inches(0.25), y + Inches(0.25), Inches(3.35), Inches(0.35),
                title, Pt(20), color, bold=True, font_name=FONT_TITLE)
    dy = y + Inches(0.85)
    for item in items:
        add_textbox(slide11, x + Inches(0.25), dy, Inches(3.35), Inches(0.28),
                    "✦ " + item, Pt(11), LIGHT_GRAY, font_name=FONT_BODY)
        dy += Inches(0.35)

# 底部提示
add_textbox(slide11, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.6),
            "💡 提示：组合使用以上功能可大幅提升开发效率。例如：用 Agent 并行探索代码库 → Memory 记录关键发现 → 定时 Cron 任务监控项目状态。",
            Pt(12), ACCENT_GOLD, font_name=FONT_BODY, alignment=PP_ALIGN.CENTER)
add_page_number(slide11, 11)

# ============================================================
# 第12页：Skills 技能系统概述
# ============================================================
slide12 = add_blank_slide()
set_bg(slide12, BG_DARK)
add_slide_title(slide12, "Skills 技能系统", "Claude Code 的专用能力模块")
add_section_indicator(slide12, ACCENT_PURPLE)

# 左侧：什么是 Skill
add_textbox(slide12, Inches(0.8), Inches(1.6), Inches(5.8), Inches(0.3),
            "什么是 Skill？", Pt(20), ACCENT_RED, bold=True, font_name=FONT_TITLE)
add_multiline_textbox(slide12, Inches(0.8), Inches(2.1), Inches(5.8), Inches(4.5), [
    ("Skill 是 Claude Code 的专用能力模块，", Pt(14), WHITE, False),
    ("每个 Skill 封装了一套特定领域的指令、流程和知识。", Pt(14), WHITE, False),
    ("", Pt(8), WHITE, False),
    ("Skill 的核心价值：", Pt(16), ACCENT_GOLD, True),
    ("", Pt(6), WHITE, False),
    ("◆ 专业化 — 让 AI 在特定领域表现更专业", Pt(13), WHITE, False),
    ("◆ 一致性 — 相同类型的任务用相同的流程处理", Pt(13), WHITE, False),
    ("◆ 效率 — 预定义流程，无需每次重复说明", Pt(13), WHITE, False),
    ("◆ 可扩展 — 支持用户自定义 Skill", Pt(13), WHITE, False),
    ("◆ 团队共享 — Skill 文件可纳入版本控制", Pt(13), WHITE, False),
    ("", Pt(8), WHITE, False),
    ("Skill 的本质是一组预先定义好的指令 prompt，", Pt(13), LIGHT_GRAY, False),
    ("被调用时加载到 Claude 的上下文中。", Pt(13), LIGHT_GRAY, False),
])

# 右侧：工作流程图
add_textbox(slide12, Inches(7.2), Inches(1.6), Inches(5.3), Inches(0.3),
            "Skill 工作原理", Pt(20), ACCENT_BLUE, bold=True, font_name=FONT_TITLE)

flow_items = [
    ("用户触发", "输入 /skill-name\n或自然语言描述需求", ACCENT_RED),
    ("系统匹配", "Claude Code 匹配\n对应的 Skill 模块", ACCENT_GOLD),
    ("加载上下文", "Skill 的指令 prompt\n被加载到对话上下文中", ACCENT_BLUE),
    ("执行任务", "Claude 按 Skill 定义的\n流程和专业标准完成任务", ACCENT_GREEN),
]

for i, (step, desc, color) in enumerate(flow_items):
    y = Inches(2.2) + i * Inches(1.15)
    add_rect(slide12, Inches(7.2), y, Inches(5.3), Inches(0.95), fill_color=BG_CARD, corner_radius=True)
    add_rect(slide12, Inches(7.22), y, Inches(0.06), Inches(0.95), fill_color=color)
    # 步骤编号
    circle = slide12.shapes.add_shape(MSO_SHAPE.OVAL, Inches(7.5), y + Inches(0.25), Inches(0.4), Inches(0.4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    p.font.name = FONT_EN
    add_textbox(slide12, Inches(8.15), y + Inches(0.1), Inches(2.0), Inches(0.28),
                step, Pt(14), WHITE, bold=True, font_name=FONT_TITLE)
    add_textbox(slide12, Inches(8.15), y + Inches(0.42), Inches(4.0), Inches(0.45),
                desc, Pt(11), LIGHT_GRAY, font_name=FONT_BODY)
add_page_number(slide12, 12)

# ============================================================
# 第13页：8 大内置 Skills
# ============================================================
slide13 = add_blank_slide()
set_bg(slide13, BG_DARK)
add_slide_title(slide13, "8 大内置 Skills", "Claude Code 自带的专业能力模块")
add_section_indicator(slide13, ACCENT_PURPLE)

skill_list = [
    ("init", "项目初始化", "自动分析项目结构，\n生成 CLAUDE.md", ACCENT_GREEN, "/init"),
    ("review", "PR 审查", "全面代码审查，\n发现 Bug 与改进点", ACCENT_BLUE, "/review"),
    ("security-review", "安全审查", "专项安全审计，\n检测漏洞与风险", ACCENT_RED, "/security-review"),
    ("simplify", "代码优化", "发现重复与低效，\n自动重构修复", ACCENT_GOLD, "/simplify"),
    ("update-config", "配置管理", "管理 settings.json，\n权限与 hooks", ACCENT_PURPLE, "/update-config"),
    ("fewer-perm", "减少弹窗", "分析历史记录，\n添加工具 allowlist", ACCENT_BLUE, "/fewer-permission-prompts"),
    ("loop", "定时循环", "定时执行任务，\n如监控部署状态", ACCENT_GREEN, "/loop 5m <task>"),
    ("claude-api", "API 开发", "Anthropic SDK 应用\n调试与优化", ACCENT_RED, "/claude-api"),
]

card_w13 = Inches(2.9)
card_h13 = Inches(2.4)
for i, (name, title, desc, color, cmd) in enumerate(skill_list):
    col = i % 4
    row = i // 4
    x = Inches(0.6) + col * (card_w13 + Inches(0.25))
    y = Inches(1.5) + row * (card_h13 + Inches(0.25))
    add_rect(slide13, x, y, card_w13, card_h13, fill_color=BG_CARD, corner_radius=True)
    add_rect(slide13, x, y, card_w13, Pt(3), fill_color=color)
    # Skill 名称
    add_textbox(slide13, x + Inches(0.2), y + Inches(0.2), card_w13 - Inches(0.4), Inches(0.3),
                name, Pt(18), color, bold=True, font_name=FONT_EN)
    # 中文标题
    add_textbox(slide13, x + Inches(0.2), y + Inches(0.55), card_w13 - Inches(0.4), Inches(0.24),
                title, Pt(12), WHITE, bold=True, font_name=FONT_TITLE)
    # 描述
    add_textbox(slide13, x + Inches(0.2), y + Inches(0.85), card_w13 - Inches(0.4), Inches(0.6),
                desc, Pt(10), LIGHT_GRAY, font_name=FONT_BODY)
    # 命令
    add_textbox(slide13, x + Inches(0.2), y + Inches(1.9), card_w13 - Inches(0.4), Inches(0.3),
                cmd, Pt(9), color, font_name=FONT_EN)
add_page_number(slide13, 13)

# ============================================================
# 第14页：如何使用 Skills
# ============================================================
slide14 = add_blank_slide()
set_bg(slide14, BG_DARK)
add_slide_title(slide14, "如何使用 Skills", "三种触发方式，灵活应对不同场景")
add_section_indicator(slide14, ACCENT_PURPLE)

# 三种方式
methods = [
    ("方式一", "斜杠命令", ACCENT_RED, Inches(0.8), Inches(1.6),
     ["在对话中直接输入 /<skill-name>",
      "最直接的调用方式",
      "示例：/init、/review、/loop 5m",
      "支持参数传递",
      "/loop 10m check deploy"]),
    ("方式二", "自然语言匹配", ACCENT_BLUE, Inches(4.8), Inches(1.6),
     ["描述任务需求",
      "系统自动匹配对应 Skill",
      "\"帮我审查这个 PR\" → review",
      "\"这段代码可以优化吗\" → simplify",
      "无需记忆命令名称"]),
    ("方式三", "后台自动调用", ACCENT_GREEN, Inches(8.8), Inches(1.6),
     ["Claude 判断任务后自动加载",
      "对用户完全透明",
      "无需手动触发",
      "匹配最合适的 Skill",
      "多 Skill 可组合使用"]),
]

for title, subtitle, color, x, y, items in methods:
    add_rect(slide14, x, y, Inches(3.7), Inches(3.8), fill_color=BG_CARD, corner_radius=True)
    add_rect(slide14, x, y, Inches(3.7), Pt(3), fill_color=color)
    add_textbox(slide14, x + Inches(0.25), y + Inches(0.15), Inches(1.2), Inches(0.24),
                title, Pt(11), color, bold=True, font_name=FONT_BODY)
    add_textbox(slide14, x + Inches(0.25), y + Inches(0.4), Inches(3.2), Inches(0.3),
                subtitle, Pt(18), WHITE, bold=True, font_name=FONT_TITLE)
    dy = y + Inches(0.95)
    for item in items:
        add_textbox(slide14, x + Inches(0.25), dy, Inches(3.2), Inches(0.24),
                    "▹ " + item, Pt(11), LIGHT_GRAY, font_name=FONT_BODY)
        dy += Inches(0.28)

# 底部提示
add_textbox(slide14, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.7),
            "💡 实际使用中，三种方式可以灵活组合。例如：手动 /init 初始化项目 → 自然语言触发 review 审查代码 → Claude 自动调用 simplify 优化代码。",
            Pt(12), ACCENT_GOLD, font_name=FONT_BODY, alignment=PP_ALIGN.CENTER)
add_page_number(slide14, 14)

# ============================================================
# 第15页：自定义 Skill
# ============================================================
slide15 = add_blank_slide()
set_bg(slide15, BG_DARK)
add_slide_title(slide15, "自定义 Skill", "打造专属你的 Claude Code 工作流")
add_section_indicator(slide15, ACCENT_PURPLE)

# 左侧：创建方式
add_textbox(slide15, Inches(0.8), Inches(1.6), Inches(5.8), Inches(0.3),
            "如何创建自定义 Skill", Pt(20), ACCENT_RED, bold=True, font_name=FONT_TITLE)
add_multiline_textbox(slide15, Inches(0.8), Inches(2.1), Inches(5.8), Inches(4.2), [
    ("1. 创建 Skill 文件", Pt(14), ACCENT_GOLD, True),
    ("在项目或用户目录的 .claude/skills/ 下", Pt(12), LIGHT_GRAY, False),
    ("创建 .md 文件，如 deploy.md", Pt(12), LIGHT_GRAY, False),
    ("", Pt(6), WHITE, False),
    ("2. 编写 Skill 指令", Pt(14), ACCENT_GOLD, True),
    ("文件内容就是 Skill 的 prompt 指令", Pt(12), LIGHT_GRAY, False),
    ("描述任务流程、检查点、输出格式", Pt(12), LIGHT_GRAY, False),
    ("可引用项目中的文件路径和命令", Pt(12), LIGHT_GRAY, False),
    ("", Pt(6), WHITE, False),
    ("3. 触发使用", Pt(14), ACCENT_GOLD, True),
    ("输入 /deploy 即可调用", Pt(12), LIGHT_GRAY, False),
    ("或通过自然语言触发", Pt(12), LIGHT_GRAY, False),
])

# 右侧：示例 Skill
add_textbox(slide15, Inches(7.2), Inches(1.6), Inches(5.3), Inches(0.3),
            "Skill 示例结构", Pt(20), ACCENT_BLUE, bold=True, font_name=FONT_TITLE)

add_rect(slide15, Inches(7.2), Inches(2.1), Inches(5.3), Inches(4.5), fill_color=BG_CARD, corner_radius=True)
add_rect(slide15, Inches(7.22), Inches(2.1), Inches(5.26), Pt(3), fill_color=ACCENT_BLUE)

example_code = [
    ("# deploy.md", Pt(11), ACCENT_RED, True),
    ("# 部署 Skill", Pt(11), LIGHT_GRAY, False),
    ("", Pt(4), WHITE, False),
    ("当用户要求部署项目时：", Pt(11), ACCENT_GOLD, True),
    ("1. 运行 npm run build", Pt(11), LIGHT_GRAY, False),
    ("2. 检查构建是否成功", Pt(11), LIGHT_GRAY, False),
    ("3. 运行对应环境的部署命令", Pt(11), LIGHT_GRAY, False),
    ("4. 确认部署后服务可访问", Pt(11), LIGHT_GRAY, False),
    ("5. 报告部署结果", Pt(11), LIGHT_GRAY, False),
    ("", Pt(4), WHITE, False),
    ("开发环境: npm run deploy:dev", Pt(11), ACCENT_GREEN, False),
    ("生产环境: npm run deploy:prod", Pt(11), ACCENT_RED, False),
    ("", Pt(4), WHITE, False),
    ("部署前务必确认分支正确。", Pt(11), ACCENT_GOLD, False),
]

dy = Inches(2.3)
for text, size, color, bold in example_code:
    add_textbox(slide15, Inches(7.5), dy, Inches(4.8), Inches(0.22),
                text, size, color, bold=bold, font_name=FONT_EN if text.startswith("#") else FONT_BODY)
    dy += Inches(0.22)

add_page_number(slide15, 15)

# ============================================================
# 第16页：最佳实践
# ============================================================
slide16 = add_blank_slide()
set_bg(slide16, BG_DARK)
add_slide_title(slide16, "最佳实践", "如何高效地与 Claude 协作")
add_section_indicator(slide16, ACCENT_GOLD)

practices = [
    ("🎯", "清晰描述需求", "具体说明技术栈、目标、约束条件，\n避免模糊的表达。好的 prompt 是成功的一半。", ACCENT_RED),
    ("📋", "善用 Plan Mode", "复杂任务先让 Claude 制定计划，\n确认方案后再开始编码，节省返工时间。", ACCENT_BLUE),
    ("🔄", "迭代式开发", "不要一次提太多需求，分步推进，\n每一步验证后再继续。", ACCENT_GREEN),
    ("📝", "编写 CLAUDE.md", "在项目中维护 CLAUDE.md 文件，\n定义项目规则、代码风格、常用指令。", ACCENT_PURPLE),
    ("🔒", "关注安全性", "不要在对话中暴露密钥和敏感信息，\n合理配置权限与审批流程。", ACCENT_RED),
    ("📊", "定期检查结果", "对 Claude 生成的代码进行验证，\n运行测试、Type Check、手动审查。", ACCENT_GOLD),
]

for i, (icon, title, desc, color) in enumerate(practices):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + col * Inches(4.1)
    y = Inches(1.6) + row * Inches(2.7)
    add_rect(slide16, x, y, Inches(3.8), Inches(2.4), fill_color=BG_CARD, corner_radius=True)
    add_rect(slide16, x, y, Inches(3.8), Pt(3), fill_color=color)
    add_textbox(slide16, x + Inches(0.25), y + Inches(0.2), Inches(0.5), Inches(0.4),
                icon, Pt(28), font_name=FONT_EN)
    add_textbox(slide16, x + Inches(0.25), y + Inches(0.75), Inches(3.3), Inches(0.3),
                title, Pt(16), WHITE, bold=True, font_name=FONT_TITLE)
    add_textbox(slide16, x + Inches(0.25), y + Inches(1.2), Inches(3.3), Inches(0.9),
                desc, Pt(11), LIGHT_GRAY, font_name=FONT_BODY)
add_page_number(slide16, 16)

# ============================================================
# 第17页：常见场景
# ============================================================
slide17 = add_blank_slide()
set_bg(slide17, BG_DARK)
add_slide_title(slide17, "常见使用场景", "实际工作中这样用 Claude Code")
add_section_indicator(slide17, ACCENT_BLUE)

scenarios = [
    ("添加新功能", ACCENT_GREEN, Inches(0.8), Inches(1.6),
     "\"在我的 Next.js 项目中\n添加用户登录功能，\n使用 JWT 认证\"",
     "→ Claude Code 分析项目结构\n→ 制定实现计划\n→ 生成 API 路由、中间件、前端组件\n→ 运行测试验证"),
    ("修复 Bug", ACCENT_RED, Inches(4.6), Inches(1.6),
     "\"用户提交表单后\n页面报 500 错误，\n帮我定位并修复\"",
     "→ 读取相关代码与错误日志\n→ 定位根本原因\n→ 提出修复方案\n→ 实施修复并验证"),
    ("学习代码库", ACCENT_PURPLE, Inches(8.4), Inches(1.6),
     "\"我刚接手这个项目，\n帮我理解整体架构\n和关键模块\"",
     "→ Agent 并行探索多个目录\n→ 绘制模块依赖关系\n→ 解释核心业务流程\n→ 标注关键入口文件"),
]

for title, color, x, y, quote, flow in scenarios:
    add_rect(slide17, x, y, Inches(3.5), Inches(2.2), fill_color=BG_CARD, corner_radius=True)
    add_rect(slide17, x, y, Inches(3.5), Pt(3), fill_color=color)
    add_textbox(slide17, x + Inches(0.2), y + Inches(0.15), Inches(3.1), Inches(0.3),
                title, Pt(18), color, bold=True, font_name=FONT_TITLE)
    # 引号区域
    add_rect(slide17, x + Inches(0.15), y + Inches(0.55), Inches(3.2), Inches(0.75),
             fill_color=RGBColor(0x20, 0x20, 0x40), corner_radius=True)
    add_textbox(slide17, x + Inches(0.3), y + Inches(0.58), Inches(2.9), Inches(0.7),
                quote, Pt(10), LIGHT_GRAY, font_name=FONT_BODY)

    # 流程
    flow_lines = flow.strip().split('\n')
    fy = y + Inches(1.5)
    for line in flow_lines:
        add_textbox(slide17, x + Inches(0.2), fy, Inches(3.1), Inches(0.22),
                    line.strip(), Pt(10), WHITE, font_name=FONT_BODY)
        fy += Inches(0.22)

# 底部：更多场景
add_textbox(slide17, Inches(0.8), Inches(4.2), Inches(11.5), Inches(0.3),
            "更多场景", Pt(16), ACCENT_GOLD, bold=True, font_name=FONT_TITLE)
more_scenarios = [
    "✧ 编写单元测试与集成测试", "✧ 数据库 Schema 设计与迁移",
    "✧ API 文档自动生成", "✧ 配置文件管理与优化",
    "✧ 依赖升级与兼容性检查", "✧ CI/CD Pipeline 配置",
]
for i, s in enumerate(more_scenarios):
    col = i % 3
    row = i // 3
    add_textbox(slide17, Inches(0.8) + col * Inches(4.1), Inches(4.7) + row * Inches(0.45),
                Inches(3.8), Inches(0.35), s, Pt(12), LIGHT_GRAY, font_name=FONT_BODY)
add_page_number(slide17, 17)

# ============================================================
# 第18页：总结
# ============================================================
slide18 = add_blank_slide()
set_bg(slide18, BG_DARK)

# 全幅装饰
add_rect(slide18, Inches(0), Inches(0), Inches(13.333), Inches(0.06), fill_color=ACCENT_RED)

# 主标题
add_textbox(slide18, Inches(1.0), Inches(1.2), Inches(11.0), Inches(0.8),
            "开始你的 Claude Code 之旅", Pt(42), WHITE, bold=True, font_name=FONT_TITLE,
            alignment=PP_ALIGN.CENTER)
add_rect(slide18, Inches(5.0), Inches(2.0), Inches(3.333), Pt(3), fill_color=ACCENT_RED)

# 关键要点
add_textbox(slide18, Inches(1.0), Inches(2.4), Inches(11.0), Inches(0.4),
            "关键要点回顾", Pt(18), ACCENT_GOLD, bold=True, font_name=FONT_TITLE,
            alignment=PP_ALIGN.CENTER)

key_points = [
    "选择适合任务的模型版本（Opus / Sonnet / Haiku）",
    "善用 Skills 系统提升专业场景效率",
    "复杂任务先用 Plan Mode 对齐方案再编码",
    "善用 Agent 并行处理提升效率",
    "通过 CLAUDE.md 和 settings.json 定制行为",
    "Memory 系统帮你积累经验、避免重复踩坑",
    "验证 Claude Code 的输出，保持人类监督",
]
for i, point in enumerate(key_points):
    y = Inches(2.95) + i * Inches(0.42)
    add_textbox(slide18, Inches(2.5), y, Inches(8.5), Inches(0.3),
                f"{i+1}.  {point}", Pt(13), WHITE, font_name=FONT_BODY)

# 资源链接
add_textbox(slide18, Inches(1.0), Inches(5.8), Inches(11.0), Inches(0.4),
            "学习资源", Pt(18), ACCENT_GOLD, bold=True, font_name=FONT_TITLE,
            alignment=PP_ALIGN.CENTER)

resources = [
    "docs.anthropic.com  — 官方文档",
    "github.com/anthropics/claude-code  — GitHub 仓库",
    "console.anthropic.com  — API 控制台",
]
for i, res in enumerate(resources):
    add_textbox(slide18, Inches(3.5), Inches(6.3) + i * Inches(0.35), Inches(6.5), Inches(0.3),
                res, Pt(12), ACCENT_BLUE, font_name=FONT_EN, alignment=PP_ALIGN.CENTER)
add_page_number(slide18, 18)

# ============================================================
# 保存
# ============================================================
output_path = os.path.expanduser(r"C:\Users\95345\Desktop\如何使用ClaudeCode.pptx")
prs.save(output_path)
print(f"PPT 已保存至: {output_path}")
print(f"共 {len(prs.slides)} 页")
