#!/usr/bin/env python3
# -*- coding: utf-8 -*-
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

DARK_BLUE   = RGBColor(0x1B, 0x2A, 0x4A)
MID_BLUE    = RGBColor(0x2E, 0x86, 0xC1)
LIGHT_BLUE  = RGBColor(0x85, 0xC1, 0xE9)
ACCENT_GOLD = RGBColor(0xF3, 0x9C, 0x12)
RED_ACCENT  = RGBColor(0xE7, 0x4C, 0x3C)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x33, 0x33, 0x33)
DARK_GRAY   = RGBColor(0x66, 0x66, 0x66)
LIGHT_BG    = RGBColor(0xF7, 0xF9, 0xFC)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

def add_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color

def add_rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def add_text(slide, l, t, w, h, text, sz=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = 'Microsoft YaHei'
    p.alignment = align
    return tb

def add_lines(slide, l, t, w, h, lines, default_sz=16, color=WHITE):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, ld in enumerate(lines):
        text = ld[0]
        sz = ld[1] if len(ld)>1 else default_sz
        bd = ld[2] if len(ld)>2 else False
        c  = ld[3] if len(ld)>3 else color
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz)
        p.font.color.rgb = c
        p.font.bold = bd
        p.font.name = 'Microsoft YaHei'
    return tb

def deco_line(slide, l, t, w, color=ACCENT_GOLD, h=Pt(4)):
    return add_rect(slide, l, t, w, h, color)

def page_num(slide, n):
    add_text(slide, Inches(12.0), Inches(7.0), Inches(1.2), Inches(0.4),
             f'{n}/10', sz=10, color=DARK_GRAY, align=PP_ALIGN.RIGHT)

def slide_title(slide, title, sub=None):
    add_rect(slide, Inches(0), Inches(0), W, Inches(1.5), DARK_BLUE)
    deco_line(slide, Inches(0.8), Inches(1.45), Inches(2.0), ACCENT_GOLD, Pt(4))
    add_text(slide, Inches(0.8), Inches(0.3), Inches(11), Inches(0.8),
             title, sz=30, color=WHITE, bold=True)
    if sub:
        add_text(slide, Inches(0.8), Inches(0.95), Inches(11), Inches(0.4),
                 sub, sz=14, color=LIGHT_BLUE)

def add_card(slide, l, t, w, h, icon, title, lines, bg=LIGHT_BG):
    add_rect(slide, l, t, w, h, bg)
    add_text(slide, l+Inches(0.2), t+Inches(0.1), Inches(0.8), Inches(0.6),
             icon, sz=28, color=MID_BLUE, bold=True)
    add_text(slide, l+Inches(0.2), t+Inches(0.6), w-Inches(0.4), Inches(0.4),
             title, sz=15, color=DARK_BLUE, bold=True)
    add_lines(slide, l+Inches(0.2), t+Inches(1.0), w-Inches(0.4), h-Inches(1.2),
              lines, default_sz=11, color=DARK_GRAY)

# ========== SLIDE 1: Cover ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_rect(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_GOLD)
add_text(slide, Inches(1.0), Inches(1.2), Inches(11), Inches(0.7),
         '邵小利青年志愿者服务队', sz=32, color=LIGHT_BLUE, bold=True)
add_text(slide, Inches(1.0), Inches(2.2), Inches(11), Inches(1.2),
         '竞聘宣传部部长', sz=54, color=WHITE, bold=True)
deco_line(slide, Inches(1.0), Inches(3.5), Inches(3.0), ACCENT_GOLD, Pt(5))
add_text(slide, Inches(1.0), Inches(3.9), Inches(8), Inches(0.6),
         '邵小利 - 个人竞选演讲', sz=22, color=LIGHT_BLUE)
add_text(slide, Inches(1.0), Inches(4.8), Inches(5), Inches(0.4),
         '2026年6月', sz=14, color=DARK_GRAY)
add_text(slide, Inches(10.0), Inches(6.6), Inches(3.0), Inches(0.4),
         '热心献社会 温暖洒人间', sz=12, color=DARK_GRAY, align=PP_ALIGN.RIGHT)
page_num(slide, 1)

# ========== SLIDE 2: Personal Profile ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_title(slide, '个人简介', 'PERSONAL PROFILE')
info = [
    ('姓    名：邵小利', 18, False, BLACK),
    ('', 8, False, BLACK),
    ('学    院：陕西中医药大学 XX学院', 18, False, BLACK),
    ('专    业：XX专业 20XX级', 18, False, BLACK),
    ('', 8, False, BLACK),
    ('现任职务：邵小利青年志愿者服务队 宣传部干事', 18, False, BLACK),
    ('', 8, False, BLACK),
    ('志愿时长：XX小时', 18, False, BLACK),
]
add_lines(slide, Inches(0.8), Inches(2.2), Inches(6.5), Inches(4.5), info, default_sz=16, color=BLACK)
add_rect(slide, Inches(8.5), Inches(2.2), Inches(4.0), Inches(4.5), LIGHT_BG)
add_text(slide, Inches(8.8), Inches(2.5), Inches(3.4), Inches(0.5),
         '自我评价', sz=20, color=DARK_BLUE, bold=True)
eval_text = [
    ('热爱志愿服务，心怀梦想', 14, True, MID_BLUE),
    ('', 6, False, BLACK),
    ('熟悉公众号运营与图文编辑', 13, False, DARK_GRAY),
    ('掌握PS、Canva、秀米等设计工具', 13, False, DARK_GRAY),
    ('良好的团队沟通与协作能力', 13, False, DARK_GRAY),
    ('对邵小利精神有深刻共鸣', 13, False, DARK_GRAY),
    ('注重细节，执行力强', 13, False, DARK_GRAY),
]
add_lines(slide, Inches(8.8), Inches(3.2), Inches(3.4), Inches(3.0), eval_text, default_sz=13, color=DARK_GRAY)
page_num(slide, 2)

# ========== SLIDE 3: Understanding the Spirit ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_title(slide, '我理解的邵小利精神', 'THE SPIRIT OF SHAO XIAOLI')
add_text(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(0.4),
         '邵小利（1960-1982），陕西中医药大学79级学生，为救落水儿童英勇牺牲，年仅22岁。',
         sz=14, color=DARK_GRAY)
add_text(slide, Inches(0.8), Inches(2.4), Inches(11.5), Inches(0.4),
         '她留下的话：生命不是私有的，不能储藏起来，如果社会和人民需要，那我就无私地贡献出来。',
         sz=14, color=RED_ACCENT, bold=True)

cards_data = [
    ('奉献', '无私助人', [
        '她让出下铺、为舍友熬药',
        '三年为同学拆洗40+床被褥',
        '把奖学金资助困难同学',
        '-> 宣传就是让奉献被看见',
    ]),
    ('担当', '舍己为人', [
        '不会游泳仍跳湖救人',
        '用双手托举孩子十几分钟',
        '22岁生命换来他人新生',
        '-> 宣传是扛起传承的责任',
    ]),
    ('精进', '追求卓越', [
        '在自己身上练针灸',
        '连续两年评为三好学生',
        '信条：宁肯在自己身上练千针',
        '-> 宣传也要精益求精',
    ]),
    ('传承', '精神不灭', [
        '1996年成立以她命名的服务队',
        '30年来一批批志愿者薪火相传',
        '校园原创话剧从未离开',
        '-> 宣传让精神接力永续',
    ]),
]
card_w = Inches(2.85)
gap = Inches(0.2)
for i, (icon, title, descs) in enumerate(cards_data):
    x = Inches(0.8) + i*(card_w+gap)
    d = [(d, 11, False, DARK_GRAY) for d in descs]
    add_card(slide, x, Inches(3.0), card_w, Inches(3.8), icon, title, d)
page_num(slide, 3)

# ========== SLIDE 4: Understanding the Role ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_title(slide, '我对宣传部工作的理解', 'MY UNDERSTANDING')
add_text(slide, Inches(0.8), Inches(2.0), Inches(11.5), Inches(0.5),
         '宣传部是服务队对外展示的窗口、对内凝聚的纽带、精神传承的扩音器',
         sz=16, color=DARK_BLUE, bold=True, align=PP_ALIGN.CENTER)

roles = [
    ('窗 口', '对外发声', [
        '让每一次志愿活动被看见',
        '传播服务队品牌形象',
        '吸引更多同学加入志愿行列',
    ]),
    ('纽 带', '对内凝聚', [
        '连接队员与服务对象',
        '记录志愿者的成长与感动',
        '增强团队归属感与荣誉感',
    ]),
    ('记 录', '留存记忆', [
        '用文字与影像定格志愿瞬间',
        '建立活动的数字档案',
        '为三十年传承留下注脚',
    ]),
    ('扩音器', '放大影响', [
        '让邵小利精神走出校园',
        '通过新媒体触达更广人群',
        '从宣传到倡导，引领志愿风尚',
    ]),
]
for i, (icon, title, descs) in enumerate(roles):
    x = Inches(0.8) + i*Inches(3.1)
    d = [(d, 12, False, DARK_GRAY) for d in descs]
    add_card(slide, x, Inches(2.8), Inches(2.8), Inches(3.8), icon, title, d)
page_num(slide, 4)

# ========== SLIDE 5: Core Competencies ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_title(slide, '我的核心能力', 'CORE COMPETENCIES')

comps = [
    ('01', '文案策划',
     '熟悉新闻稿、活动推文、倡议书等多种文体\n参与服务队多项活动的宣传方案策划\n注重文字的温度与感染力',
     ACCENT_GOLD),
    ('02', '视觉设计',
     '熟练使用PS、Canva、秀米等设计工具\n累计设计海报30+张、推文封面80+张\n注重志愿主题的视觉表达',
     MID_BLUE),
    ('03', '新媒体运营',
     '熟悉公众号、抖音等平台规则\n掌握排版美学与数据分析\n善于用年轻人喜欢的方式讲志愿故事',
     RED_ACCENT),
    ('04', '摄影记录',
     '活动全程跟拍摄影经验\n注重捕捉真实、感人的志愿瞬间\n照片即宣传素材，也是服务队的历史',
     DARK_BLUE),
]
for i, (num, title, desc, accent) in enumerate(comps):
    y = Inches(2.0) + i*Inches(1.35)
    add_text(slide, Inches(0.8), y+Inches(0.1), Inches(0.7), Inches(0.6),
             num, sz=32, color=accent, bold=True)
    add_rect(slide, Inches(1.5), y+Inches(0.2), Pt(4), Inches(0.6), accent)
    add_text(slide, Inches(1.8), y, Inches(4), Inches(0.5),
             title, sz=20, color=DARK_BLUE, bold=True)
    add_text(slide, Inches(1.8), y+Inches(0.5), Inches(10), Inches(0.8),
             desc, sz=12, color=DARK_GRAY)
page_num(slide, 5)

# ========== SLIDE 6: Achievements ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_title(slide, '过往工作成果', 'ACHIEVEMENTS')

projects = [
    ('XX志愿活动宣传', '阅读量2000+  转发150+', [
        '主导整体宣传策划，统筹线上线下推广',
        '设计系列海报6张，编写预热/回顾推文10篇',
        '活动当天负责全程摄影记录',
    ]),
    ('服务队招新宣传', '阅读量3000+  新增报名80+', [
        '策划招新系列推文，设计招新海报与展板',
        '抖音招新短视频累计播放量1万+',
        '创新线上报名+线下宣讲结合模式',
    ]),
    ('日常公众号运营', '粉丝增长40%  推文产出30+篇', [
        '负责公众号每周推文的撰写与排版',
        '建立素材库，提高内容生产效率',
        '优化视觉风格，提升公众号辨识度',
    ]),
]
for i, (title, metric, descs) in enumerate(projects):
    y = Inches(2.0) + i*Inches(1.8)
    add_rect(slide, Inches(0.8), y, Pt(4), Inches(1.2), ACCENT_GOLD)
    add_text(slide, Inches(1.1), y, Inches(5), Inches(0.4),
             title, sz=18, color=DARK_BLUE, bold=True)
    add_text(slide, Inches(1.1), y+Inches(0.4), Inches(5), Inches(0.3),
             metric, sz=12, color=MID_BLUE)
    d = [(d, 12, False, DARK_GRAY) for d in descs]
    add_lines(slide, Inches(6.5), y+Inches(0.1), Inches(6), Inches(1.2), d, default_sz=12, color=DARK_GRAY)
page_num(slide, 6)

# ========== SLIDE 7: Work Plan ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_title(slide, '工作规划', 'WORK PLAN')

phases = [
    ('短期（1-3个月）', ACCENT_GOLD, [
        '建立标准化宣传工作流程与排期制度',
        '完成部门成员技能摸底与分工优化',
        '统一公众号视觉风格，提高辨识度',
        '完善活动素材库，提高内容复用率',
    ]),
    ('中期（3-6个月）', MID_BLUE, [
        '拓展抖音/视频号等短视频宣传渠道',
        '策划1-2个志愿服务品牌宣传项目',
        '公众号粉丝增长50%，阅读量翻倍',
        '加强与校内外媒体的合作对接',
    ]),
    ('长期（6-12个月）', DARK_BLUE, [
        '打造邵小利服务队校园IP，扩大影响力',
        '建立跨部门宣传协作机制',
        '培养3-5名宣传骨干，形成人才梯队',
        '总结方法论，为下一届留下可复用经验文档',
    ]),
]
for i, (phase, color, items) in enumerate(phases):
    y = Inches(2.0) + i*Inches(1.8)
    add_rect(slide, Inches(0.8), y+Inches(0.05), Inches(0.12), Inches(0.4), color)
    add_text(slide, Inches(1.2), y, Inches(5), Inches(0.4),
             phase, sz=17, color=color, bold=True)
    items_text = '\n'.join(['● '+item for item in items])
    add_text(slide, Inches(1.2), y+Inches(0.5), Inches(11), Inches(1.2),
             items_text, sz=12, color=DARK_GRAY)
page_num(slide, 7)

# ========== SLIDE 8: Innovation ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_title(slide, '创新思路', 'INNOVATIVE IDEAS')

ideas = [
    ('故事化宣传',
     '从活动通知到志愿故事\n用真实的人物和故事打动人心\n让每一次宣传都是一次精神传递'),
    ('短视频矩阵',
     '抖音+视频号双平台运营\n15-60秒记录志愿精彩瞬间\n用年轻人喜欢的方式传播正能量'),
    ('互动升级',
     '从单向推送到双向互动\n话题征集/志愿者专访/UGC内容\n让每位队员都成为内容的创造者'),
    ('品牌化建设',
     '打造服务队专属视觉体系\n标志性海报风格+统一排版\n让邵小利服务队一眼可辨'),
]
for i, (title, desc) in enumerate(ideas):
    x = Inches(0.8) + i*Inches(3.1)
    add_rect(slide, x, Inches(2.0), Inches(2.8), Inches(0.06), MID_BLUE)
    add_text(slide, x, Inches(2.3), Inches(2.8), Inches(0.5),
             title, sz=20, color=DARK_BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x, Inches(2.9), Inches(2.8), Inches(3.0),
             desc, sz=12, color=DARK_GRAY, align=PP_ALIGN.CENTER)
page_num(slide, 8)

# ========== SLIDE 9: My Commitment ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_title(slide, '我的承诺', 'MY COMMITMENT')

commitments = [
    ('传承精神',
     '以邵小利精神为指引，用心做好每一次宣传\n让热心献社会，温暖洒人间的宗旨\n通过宣传传递给更多人'),
    ('服务团队',
     '做好幕后工作者，把聚光灯留给一线志愿者\n建立宣传素材即提即用的高效响应机制\n让每个部门的活动都能得到最好的展示'),
    ('不断成长',
     '关注高校宣传新趋势，持续学习新技能\n定期复盘，用数据驱动内容优化\n带领宣传部成为学习型团队'),
    ('用心记录',
     '不让任何一个志愿故事被遗忘\n建立服务队年度大事记与影像年鉴\n为三十年历史续写新的篇章'),
]
colors = [ACCENT_GOLD, MID_BLUE, RED_ACCENT, DARK_BLUE]
for i, (title, desc) in enumerate(commitments):
    y = Inches(2.0) + i*Inches(1.35)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y+Inches(0.15),
                                  Inches(0.4), Inches(0.4))
    dot.fill.solid()
    dot.fill.fore_color.rgb = colors[i]
    dot.line.fill.background()
    add_text(slide, Inches(0.8), y+Inches(0.15), Inches(0.4), Inches(0.4),
             str(i+1), sz=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.5), y, Inches(4), Inches(0.5),
             title, sz=20, color=DARK_BLUE, bold=True)
    add_text(slide, Inches(1.5), y+Inches(0.5), Inches(11), Inches(0.8),
             desc, sz=12, color=DARK_GRAY)
page_num(slide, 9)

# ========== SLIDE 10: Closing ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_rect(slide, Inches(0), Inches(0), Inches(0.15), H, ACCENT_GOLD)
add_text(slide, Inches(1.0), Inches(2.0), Inches(11), Inches(0.8),
         '感谢聆听', sz=48, color=WHITE, bold=True)
deco_line(slide, Inches(1.0), Inches(3.0), Inches(2.0), ACCENT_GOLD, Pt(4))
add_text(slide, Inches(1.0), Inches(3.4), Inches(11), Inches(0.6),
         '热心献社会，温暖洒人间', sz=22, color=LIGHT_BLUE)
add_text(slide, Inches(1.0), Inches(4.3), Inches(8), Inches(0.5),
         '邵小利 - 宣传部部长竞聘', sz=16, color=DARK_GRAY)
add_text(slide, Inches(1.0), Inches(4.9), Inches(8), Inches(0.5),
         '期待与大家同行，让志愿精神薪火相传', sz=14, color=DARK_GRAY)
page_num(slide, 10)

# Save
output = 'C:/Users/95345/Desktop/邵小利_宣传部部长竞聘.pptx'
prs.save(output)
print(f'Done! Saved to: {output}')
print(f'Total slides: {len(prs.slides)}')
